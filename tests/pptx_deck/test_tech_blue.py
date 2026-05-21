# tests/pptx_deck/test_tech_blue.py
"""tech_blue 主题 11 layout light test：验证每个 layout 创建后 prs.slides 增加 1。"""
import os, sys
ROOT = os.path.join(os.path.dirname(__file__), "../../skills")
sys.path.insert(0, os.path.join(ROOT, "pptx"))
sys.path.insert(0, os.path.join(ROOT, "pptx-deck"))

from pptx import Presentation
from pptx.util import Inches
from themes import tech_blue as T


def _new():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def test_make_cover():
    prs = _new()
    T.make_cover(prs, "主标题", "副标题")
    assert len(prs.slides) == 1

def test_make_toc():
    prs = _new()
    T.make_toc(prs, sections=["背景", "范围", "流程", "保障", "节奏"])
    assert len(prs.slides) == 1

def test_make_section_divider():
    prs = _new()
    T.make_section_divider(prs, 1, "第一章")
    assert len(prs.slides) == 1

def test_make_single_focus():
    prs = _new()
    T.make_single_focus(prs, big_text="一句话", big_number="80%", explanation="解释")
    assert len(prs.slides) == 1

def test_make_two_col_compare():
    prs = _new()
    T.make_two_col_compare(prs, "现状", "现状描述", "目标", "目标描述")
    assert len(prs.slides) == 1

def test_make_three_col_cards():
    prs = _new()
    T.make_three_col_cards(prs, cards=[
        {"title": "卡1", "body": "正文1"},
        {"title": "卡2", "body": "正文2"},
        {"title": "卡3", "body": "正文3"},
    ])
    assert len(prs.slides) == 1

def test_make_bullet_list():
    prs = _new()
    T.make_bullet_list(prs, "标题", items=["要点1", "要点2", "要点3", "要点4", "要点5"])
    assert len(prs.slides) == 1

def test_make_table():
    prs = _new()
    T.make_table(prs, "表格标题",
                 headers=["A", "B", "C"],
                 rows=[["1", "2", "3"], ["4", "5", "6"]])
    assert len(prs.slides) == 1

def test_make_pic_text(tmp_path):
    from PIL import Image
    img_path = tmp_path / "blank.png"
    Image.new("RGB", (10, 10), "white").save(str(img_path))
    prs = _new()
    T.make_pic_text(prs, "标题", str(img_path),
                    points=[{"title": "点1", "body": "正文1"},
                            {"title": "点2", "body": "正文2"}])
    assert len(prs.slides) == 1

def test_make_summary():
    prs = _new()
    T.make_summary(prs, conclusions=["结论 1", "结论 2", "结论 3"])
    assert len(prs.slides) == 1

def test_make_closing():
    prs = _new()
    T.make_closing(prs, subtitle="联系邮箱：x@y.com")
    assert len(prs.slides) == 1

def test_font_default_is_microsoft_yahei():
    assert T.FONT_HEADER == "Microsoft YaHei"
    assert T.FONT_BODY == "Microsoft YaHei"
