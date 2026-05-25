"""build.py 单元测试。"""
import json
from pathlib import Path

import pytest

from pptx import Presentation as _Pres

from build import (
    FOOTERED_LAYOUTS,
    _extract_design_tokens,
    build_deck,
    load_plan,
    load_theme,
)
from themes import tech_blue


def _write_plan(tmp_path, data):
    p = tmp_path / "plan.json"
    p.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
    return p


def test_load_plan_valid(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./out.pptx",
        "slides": [{"layout": "cover", "title": "T", "subtitle": "S"}]})
    plan = load_plan(p)
    assert plan["theme"] == "tech_blue"
    assert plan["_plan_dir"] == str(tmp_path.resolve())


def test_load_plan_missing_field_raises(tmp_path):
    p = _write_plan(tmp_path, {"theme": "tech_blue", "slides": []})
    with pytest.raises(ValueError) as e:
        load_plan(p)
    assert "output" in str(e.value)


def test_load_plan_slide_missing_layout_raises(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./o.pptx",
        "slides": [{"title": "无 layout"}]})
    with pytest.raises(ValueError) as e:
        load_plan(p)
    assert "第 1 页" in str(e.value)


def test_load_theme_tech_blue():
    assert load_theme("tech_blue") is tech_blue


def test_load_theme_unknown_raises():
    with pytest.raises(ValueError):
        load_theme("nope")


def test_load_theme_pptx_missing_raises():
    with pytest.raises(FileNotFoundError):
        load_theme("/tmp/does_not_exist_xyz.pptx")


def test_load_theme_short_name_resolves_via_plan_templates(tmp_path):
    """deck 工作目录下 templates/<name>.pptx 优先生效"""
    from pptx import Presentation as _P
    templates_dir = tmp_path / "templates"
    templates_dir.mkdir()
    _P().save(str(templates_dir / "company_a.pptx"))
    theme = load_theme("company_a", plan_dir=str(tmp_path))
    # 派生 theme 模块至少有 make_cover(从 tech_blue 复制 + 改 token)
    assert hasattr(theme, "make_cover")
    assert hasattr(theme, "FONT_HEADER")
    assert hasattr(theme, "PRIMARY")


def test_load_theme_short_name_not_found_lists_available(tmp_path):
    """短名查不到时报错应当列出 templates/ 现有名字 + 提示"""
    with pytest.raises(ValueError) as e:
        load_theme("nonexistent_xyz", plan_dir=str(tmp_path))
    msg = str(e.value)
    assert "nonexistent_xyz" in msg
    assert "tech_blue" in msg               # 提示内置可选
    assert "templates/" in msg              # 提示查 templates/


def test_load_theme_per_plan_overrides_repo_templates(tmp_path):
    """plan_dir 下的 templates 优先于 repo templates"""
    from pptx import Presentation as _P
    plan_templates = tmp_path / "templates"
    plan_templates.mkdir()
    # 用一个明显不一样的内容验证用的是 plan 的而非 repo 的
    _P().save(str(plan_templates / "override_test.pptx"))
    # plan_dir 命中
    theme = load_theme("override_test", plan_dir=str(tmp_path))
    assert hasattr(theme, "make_cover")
    # 不传 plan_dir 时,仅 repo templates 找,这个 override_test 没在 repo,应 raise
    with pytest.raises(ValueError):
        load_theme("override_test")


# ----- 扩展 token 提取测试(Tier 1 模板摄入)-----

def test_extract_extended_tokens_returns_dict(tmp_path):
    """_extract_design_tokens 应返回 dict,即使模板很简单也不报错"""
    from pptx import Presentation as _P
    from build import _extract_design_tokens
    p = tmp_path / "minimal.pptx"
    _P().save(str(p))
    tokens = _extract_design_tokens(str(p))
    # 至少不抛异常;空 pptx 可能没有 accent1 / ea typeface,返回空 dict 或部分
    assert isinstance(tokens, dict)


def test_extract_template_cli_l1_l2(tmp_path):
    """extract_template.py CLI 跑通(L1 媒体提取 + L2 token,跳过 probe 加速)"""
    import subprocess
    from pptx import Presentation as _P
    from pathlib import Path as _Pa
    # 1. 写一个 .pptx 模板到 templates/
    tmpl_dir = tmp_path / "templates"
    tmpl_dir.mkdir()
    pptx = tmpl_dir / "minimal_template.pptx"
    _P().save(str(pptx))
    # 2. 跑 CLI
    script = _Pa(__file__).resolve().parent.parent.parent / "skills" / "pptx-deck" / "extract_template.py"
    result = subprocess.run(
        ["python3", str(script), str(pptx),
         "--no-probe", "--working-dir", str(tmp_path)],
        capture_output=True, text=True)
    assert result.returncode == 0, f"CLI failed: {result.stderr}"
    # 3. yaml 应被写
    yaml_path = tmpl_dir / "minimal_template.yaml"
    assert yaml_path.exists()
    # 4. yaml 应包含 extracted 段
    import yaml as _y
    data = _y.safe_load(yaml_path.read_text(encoding="utf-8"))
    assert "extracted" in data
    assert "source_pptx" in data["extracted"]
    assert "tokens" in data["extracted"]


def test_extract_template_preserves_user_yaml_fields(tmp_path):
    """已有 yaml 的用户字段(desc/notes/owner)应保留,不被 extract 覆盖"""
    import subprocess
    from pptx import Presentation as _P
    from pathlib import Path as _Pa
    import yaml as _y
    tmpl_dir = tmp_path / "templates"
    tmpl_dir.mkdir()
    pptx = tmpl_dir / "test_preserve.pptx"
    _P().save(str(pptx))
    # 预先写一份 yaml 含用户字段
    yaml_path = tmpl_dir / "test_preserve.yaml"
    yaml_path.write_text(
        "name: 我的模板\ndesc: 不能覆盖\nowner: alice\nnotes: 保留我\n",
        encoding="utf-8")
    # 跑 CLI
    script = _Pa(__file__).resolve().parent.parent.parent / "skills" / "pptx-deck" / "extract_template.py"
    result = subprocess.run(
        ["python3", str(script), str(pptx), "--no-probe", "--working-dir", str(tmp_path)],
        capture_output=True, text=True)
    assert result.returncode == 0
    # 验证用户字段被保留
    data = _y.safe_load(yaml_path.read_text(encoding="utf-8"))
    assert data["name"] == "我的模板"
    assert data["desc"] == "不能覆盖"
    assert data["owner"] == "alice"
    assert data["notes"] == "保留我"
    # 同时 extracted 字段被加上
    assert "extracted" in data


def test_build_deck_produces_pptx(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./deck.pptx",
        "slides": [
            {"layout": "cover", "title": "标题", "subtitle": "副标题"},
            {"layout": "bullet_list", "title": "要点", "items": ["a", "b"]}]})
    plan = load_plan(p)
    out = build_deck(plan)
    assert out.exists()
    assert out.stat().st_size > 0
    assert out.parent == tmp_path.resolve()


def test_build_deck_unknown_layout_raises(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx",
        "slides": [{"layout": "nonexistent_xyz"}]})
    plan = load_plan(p)
    with pytest.raises(ValueError) as e:
        build_deck(plan)
    assert "未知 layout" in str(e.value)


def test_build_deck_bad_field_raises_with_page_number(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx",
        "slides": [{"layout": "cover", "wrong_field": "x"}]})
    plan = load_plan(p)
    with pytest.raises(ValueError) as e:
        build_deck(plan)
    assert "第 1 页" in str(e.value)


# ----- footer / 页码集成测试 -----

def _build_and_open(tmp_path, slides_def):
    """构建一个 mini deck,返回打开后的 Presentation 对象供 inspect。"""
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx", "slides": slides_def})
    plan = load_plan(p)
    out = build_deck(plan)
    return _Pres(str(out))


def test_footer_added_only_to_footered_layouts(tmp_path):
    """cover / section_divider / closing 不带页脚;其他 8 种 layout 带。"""
    prs = _build_and_open(tmp_path, [
        {"layout": "cover", "title": "T", "subtitle": "S"},
        {"layout": "section_divider", "num": 1, "title": "第一章"},
        {"layout": "bullet_list", "title": "要点", "items": ["a", "b"]},
        {"layout": "closing", "subtitle": ""},
    ])
    # 找带 "/" 的 textbox(页码标记)
    def has_page_num(slide):
        for sh in slide.shapes:
            if sh.has_text_frame and "/" in sh.text_frame.text:
                # 排除大段含"/"的正文:页码一定是 "N / M" 整段
                t = sh.text_frame.text.strip()
                if len(t) <= 10 and t.replace(" ", "").replace("/", "").isdigit():
                    return True
        return False
    assert not has_page_num(prs.slides[0]), "cover 不应有页码"
    assert not has_page_num(prs.slides[1]), "section_divider 不应有页码"
    assert has_page_num(prs.slides[2]), "bullet_list 应有页码"
    assert not has_page_num(prs.slides[3]), "closing 不应有页码"


def test_footer_page_num_uses_content_page_index_not_absolute(tmp_path):
    """N / TOTAL 的 N 是 content page 序号(跳过 cover/divider/closing),
    TOTAL 是 content page 总数。"""
    prs = _build_and_open(tmp_path, [
        {"layout": "cover", "title": "T", "subtitle": "S"},
        {"layout": "section_divider", "num": 1, "title": "一"},
        {"layout": "cards", "title": "A",
         "cards": [{"title": "x", "body": "y"}]},
        {"layout": "section_divider", "num": 2, "title": "二"},
        {"layout": "cards", "title": "B",
         "cards": [{"title": "x", "body": "y"}]},
        {"layout": "closing", "subtitle": ""},
    ])
    # 2 个 content page(cards × 2)
    page_nums = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if "/" in t and len(t) <= 10:
                    page_nums.append(t)
    assert page_nums == ["1 / 2", "2 / 2"]


def test_footered_layouts_constant_matches_spec():
    """常量必须覆盖规范要求的 10 种核心 layout + visual pattern 衍生 layout(除 cover/section_divider/closing)。"""
    assert FOOTERED_LAYOUTS == frozenset({
        "toc", "single_focus", "compare", "compare_pk", "matrix_2x2", "cards",
        "bullet_list", "table", "pic_text", "summary",
        # visual patterns 衍生 layout(template_training theme 引入)
        "timeline_band_3", "tri_pyramid_4sub_3", "cards_flag_3",
    })


def test_source_field_renders_citation(tmp_path):
    """deck_plan slide 含 'source' 字段时,build.py 自动加 Source: 引文。"""
    prs = _build_and_open(tmp_path, [
        {"layout": "table", "title": "数据", "headers": ["A", "B"],
         "rows": [["1", "2"]],
         "source": "公司财报 2025 Q4"},
    ])
    texts = " | ".join(
        sh.text_frame.text for sh in prs.slides[0].shapes
        if sh.has_text_frame and sh.text_frame.text
    )
    assert "Source: 公司财报 2025 Q4" in texts


def test_source_field_does_not_break_make_table_signature(tmp_path):
    """'source' 是 cross-cutting 字段,不应作为 make_table 的 kwarg。"""
    # 如果 build.py 没正确 pop "source",会触发 TypeError
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx",
        "slides": [{"layout": "table", "title": "T",
                    "headers": ["A"], "rows": [["1"]],
                    "source": "Test"}]})
    plan = load_plan(p)
    out = build_deck(plan)
    assert out.exists()


def test_footer_meta_propagates_to_all_content_pages(tmp_path):
    """plan.footer_meta 应在每个内容页 footer 左侧显示。"""
    prs = _build_and_open(tmp_path, [
        {"layout": "cover", "title": "T", "subtitle": "S"},
        {"layout": "cards", "title": "A",
         "cards": [{"title": "x", "body": "y"}]},
        {"layout": "bullet_list", "title": "B", "items": ["1", "2"]},
        {"layout": "closing", "subtitle": ""},
    ])
    # 在 deck_plan 顶层加 footer_meta —— 需要走 _write_plan 写完整 plan
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d2.pptx",
        "footer_meta": {
            "classification": "INTERNAL",
            "project": "Project Atlas",
            "version": "v1.0",
        },
        "slides": [
            {"layout": "cover", "title": "T", "subtitle": "S"},
            {"layout": "cards", "title": "A",
             "cards": [{"title": "x", "body": "y"}]},
            {"layout": "bullet_list", "title": "B", "items": ["1", "2"]},
            {"layout": "closing", "subtitle": ""},
        ],
    })
    plan = load_plan(p)
    out = build_deck(plan)
    prs2 = _Pres(str(out))
    # cards 和 bullet_list 应有 footer_meta;cover 和 closing 没有
    def has_meta(slide):
        for sh in slide.shapes:
            if sh.has_text_frame and "INTERNAL · Project Atlas · v1.0" in sh.text_frame.text:
                return True
        return False
    assert not has_meta(prs2.slides[0]), "cover 不应有 footer_meta"
    assert has_meta(prs2.slides[1]), "cards 应有 footer_meta"
    assert has_meta(prs2.slides[2]), "bullet_list 应有 footer_meta"
    assert not has_meta(prs2.slides[3]), "closing 不应有 footer_meta"
