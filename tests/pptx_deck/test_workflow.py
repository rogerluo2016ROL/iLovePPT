"""workflow.py 单元测试 — 覆盖错误路径与核心 dispatch 逻辑。"""
import tempfile
from pathlib import Path

import pytest
import yaml

from workflow import (
    parse_brief,
    load_theme,
    estimate_page_count,
    generate_outline,
    generate_slide,
    fix_slide,
    REQUIRED,
)
from themes import tech_blue


# ----- parse_brief -----

def test_parse_brief_loads_valid_yaml(tmp_path):
    p = tmp_path / "b.yaml"
    p.write_text(yaml.safe_dump({
        "title": "T", "outline": ["a"], "theme": "tech_blue",
        "output": "/tmp/x.pptx"
    }))
    b = parse_brief(str(p))
    assert b["title"] == "T"
    assert b["outline"] == ["a"]
    # 默认字段被填充
    assert b["subtitle"] == ""
    assert b["page_count_target"] is None
    assert b["key_points"] == []
    assert b["reference_pptx"] is None


def test_parse_brief_missing_required_raises(tmp_path):
    p = tmp_path / "b.yaml"
    p.write_text(yaml.safe_dump({"title": "T"}))
    with pytest.raises(ValueError) as exc:
        parse_brief(str(p))
    assert "缺字段" in str(exc.value)
    # 错误提示应包含三个缺失字段名
    assert "outline" in str(exc.value)
    assert "theme" in str(exc.value)
    assert "output" in str(exc.value)


def test_parse_brief_required_set_correct():
    assert REQUIRED == {"title", "outline", "theme", "output"}


# ----- load_theme -----

def test_load_theme_tech_blue_returns_module():
    t = load_theme("tech_blue")
    assert t is tech_blue


def test_load_theme_pptx_file_not_exists_raises():
    with pytest.raises(FileNotFoundError):
        load_theme("/tmp/this_file_does_not_exist_iloveppt.pptx")


def test_load_theme_unknown_raises_value_error():
    with pytest.raises(ValueError) as exc:
        load_theme("non_existent_theme")
    assert "未知 theme" in str(exc.value)


# ----- estimate_page_count -----

def test_estimate_page_count_uses_target_if_set():
    brief = {"outline": ["a", "b", "c"], "page_count_target": 20}
    assert estimate_page_count(brief) == 20


def test_estimate_page_count_formula_when_no_target():
    # 4 sections × 1.5 + 4 (cover/toc/summary/closing) = 10
    brief = {"outline": ["a", "b", "c", "d"], "page_count_target": None}
    assert estimate_page_count(brief) == 10


# ----- generate_outline -----

def test_generate_outline_produces_expected_structure():
    brief = {
        "title": "T",
        "subtitle": "S",
        "outline": ["sec1", "sec2"],
        "key_points": ["kp1", "kp2"],
        "theme": "tech_blue",
    }
    specs = generate_outline(brief)
    # cover + toc + (section_divider + bullet_list) × 2 + summary + closing = 8
    assert len(specs) == 8
    assert specs[0]["layout"] == "cover"
    assert specs[1]["layout"] == "toc"
    assert specs[2]["layout"] == "section_divider"
    assert specs[2]["num"] == 1
    assert specs[3]["layout"] == "bullet_list"
    assert specs[-2]["layout"] == "summary"
    assert specs[-1]["layout"] == "closing"


def test_generate_outline_empty_key_points_uses_default():
    brief = {
        "title": "T", "subtitle": "",
        "outline": ["s1"], "key_points": [],
    }
    specs = generate_outline(brief)
    bullet_spec = next(s for s in specs if s["layout"] == "bullet_list")
    # 默认 items 应非空（fallback "{sec} 要点 1/2"）
    assert len(bullet_spec["items"]) >= 2


# ----- generate_slide dispatch -----

from pptx import Presentation
from pptx.util import Inches


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def test_generate_slide_dispatches_to_make_cover():
    prs = _new_prs()
    spec = {"layout": "cover", "title": "T", "subtitle": "S"}
    generate_slide(prs, spec, tech_blue)
    assert len(prs.slides) == 1


def test_generate_slide_dispatches_to_make_bullet_list():
    prs = _new_prs()
    spec = {"layout": "bullet_list", "title": "T", "items": ["a", "b"]}
    generate_slide(prs, spec, tech_blue)
    assert len(prs.slides) == 1


def test_generate_slide_unknown_layout_raises_attribute_error():
    prs = _new_prs()
    spec = {"layout": "nonexistent_layout"}
    with pytest.raises(AttributeError):
        generate_slide(prs, spec, tech_blue)


# ----- fix_slide -----

def test_fix_slide_no_issues_returns_slide_unchanged():
    prs = _new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    result = fix_slide(slide, [])
    assert result is slide


def test_fix_slide_fontsize_too_large_reduces_fonts():
    from pptx.util import Inches, Pt
    prs = _new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "test"
    run.font.size = Pt(40)
    issues = [{"issue": "test", "suggested_fix": "字号过大"}]
    fix_slide(slide, issues)
    assert run.font.size == Pt(32)  # 40 × 0.8 = 32


def test_fix_slide_margin_not_zeroed_fixes_margins():
    from pptx.util import Inches, Emu
    prs = _new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    issues = [{"issue": "test", "suggested_fix": "margin 未归零"}]
    fix_slide(slide, issues)
    assert tb.text_frame.margin_left == Emu(0)


def test_fix_slide_unknown_fix_does_not_modify_shape_count():
    prs = _new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    n_shapes = len(slide.shapes)
    fix_slide(slide, [{"issue": "test", "suggested_fix": "未知修复策略 xyz"}])
    assert len(slide.shapes) == n_shapes


# ----- _ingest_template / load_theme with .pptx -----

def test_load_theme_pptx_ingest_uses_minimal_deck(tmp_path):
    """用 minimal_deck.py 生成的 .pptx 作为 ingest 输入,验证 load_theme 返回带 make_* 的 module。"""
    import importlib.util
    # minimal_deck 位于 skills/pptx/examples/,通过绝对路径加载
    here = Path(__file__).parent
    md_path = here.parent.parent / "skills" / "pptx" / "examples" / "minimal_deck.py"
    spec = importlib.util.spec_from_file_location("minimal_deck", md_path)
    md = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(md)

    out = str(tmp_path / "ingest_test.pptx")
    md.main(out=out)

    theme = load_theme(out)
    for name in [
        "make_cover", "make_toc", "make_section_divider",
        "make_single_focus", "make_two_col_compare",
        "make_three_col_cards", "make_bullet_list", "make_table",
        "make_pic_text", "make_summary", "make_closing",
    ]:
        assert hasattr(theme, name), f"ingested theme 缺少 {name}"
