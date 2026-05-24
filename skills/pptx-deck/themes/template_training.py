"""template_training 主题 —— iSlide 商业模板的全 13 layout 复刻(培训 / 内训场景)

# 设计原则
- 双主色:橙红 #EF5938(强调 / 章节号 / accent 条) + 深蓝 #0B2A4A(标题 / VS 圆)
- 标志性元素:TEAM 四人扁平插画(_assets/template_template_training/team_illustration.png)
- 5 个结构性 layout(cover/toc/section_divider/cards/closing)按 native render 重设计
- 8 个内容 layout 沿用 tech_blue 结构 + 切换到 template_training 调色板

# 复刻范围 vs native template
- Tier 1 重写:cover/toc/section_divider/cards/closing(含 TEAM 插画 + LOGO 位)
- Tier 2 切色:single_focus/compare/compare_pk/matrix_2x2/bullet_list/table/pic_text/summary
- 不复刻:SmartArt 装饰图(辐射 / 4 象限带剪影)、iSlide footer 水印

# 适用场景
- 企业内训 / 课程开训 / 团队培训 / 复盘汇报 / 入职引导
- 高对比 + 暖色 → 抓注意力强
"""
import sys
from pathlib import Path
from typing import Any

# Fallback for direct import outside pytest
_helpers_path = str(Path(__file__).parent.parent.parent / "pptx")
if _helpers_path not in sys.path:
    sys.path.insert(0, _helpers_path)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as _Pres
from pptx.slide import Slide
from pptx.util import Emu, Inches, Length, Pt

import helpers as H
import layout as L


# ===== 字体 / 色板 =====
FONT_HEADER = H.FONT_CN
FONT_BODY   = H.FONT_CN
FONT_NUM    = H.FONT_NUM

PRIMARY      = RGBColor(0xEF, 0x59, 0x38)   # 橙红(模板 accent1)
PRIMARY_DEEP = RGBColor(0x0B, 0x2A, 0x4A)   # 深蓝(继承 tech_blue BRAND_DARK)
PRIMARY_TINT = RGBColor(0xFD, 0xE5, 0xDE)   # 浅橙底
LIGHT_BLUE   = RGBColor(0xC8, 0xE4, 0xF8)   # 浅蓝带(section_divider 下)
ACCENT       = RGBColor(0x08, 0x59, 0x86)   # 深蓝(模板 accent2)

# 资产路径(指仓库根的 _assets/template_template_training/)
_REPO_ROOT = Path(__file__).resolve().parent.parent.parent.parent
TEAM_ILLUSTRATION = _REPO_ROOT / "_assets" / "template_template_training" / "team_illustration.png"


# ============================================================================
# 内部 helpers
# ============================================================================

def _blank_slide(prs: _Pres) -> Slide:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title(slide: Slide, text: str, *, y: Length = Inches(0.6),
               size: int = 32, color: RGBColor = PRIMARY_DEEP) -> Any:
    """页面标题 — 跟 tech_blue 一致 32pt 起。"""
    box = slide.shapes.add_textbox(Inches(0.55), y, Inches(12.2), Inches(0.9))
    tf = box.text_frame
    H.fix_textbox_margins(tf)
    r = tf.paragraphs[0].add_run()
    r.text = text
    H.set_font(r, name=FONT_HEADER, size=size, bold=True, color=color)
    return box


def _text_box(slide: Slide, box: "L.Box", text: str, *, size: int,
              bold: bool = False, color=None, align=PP_ALIGN.LEFT,
              font=None) -> None:
    """tech_blue 风格:在一个 L.Box 内放文字。"""
    if color is None:
        color = H.GRAY_900
    if font is None:
        font = FONT_HEADER
    tb = slide.shapes.add_textbox(box.x, box.y, box.w, box.h)
    H.fix_textbox_margins(tb.text_frame)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    H.set_font(r, name=font, size=size, bold=bold, color=color)


def _text_xy(slide: Slide, x: Length, y: Length, w: Length, h: Length,
             text: str, *, size: int, bold: bool = False, color=None,
             align=PP_ALIGN.LEFT, font=None) -> None:
    """template_a 风格:绝对坐标放文字。"""
    if color is None:
        color = H.GRAY_900
    if font is None:
        font = FONT_HEADER
    tb = slide.shapes.add_textbox(x, y, w, h)
    H.fix_textbox_margins(tb.text_frame)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    H.set_font(r, name=font, size=size, bold=bold, color=color)


def _logo_placeholder(slide: Slide, x: Length, y: Length) -> None:
    """LOGO 占位框 — 浅灰描边小框 + 'LOGO' 文字"""
    w, h = Inches(0.9), Inches(0.4)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.background()
    rect.line.color.rgb = H.GRAY_300
    rect.line.width = Pt(0.5)
    tf = rect.text_frame
    H.fix_textbox_margins(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "LOGO"
    H.set_font(r, name=FONT_HEADER, size=10, color=H.GRAY_500)


# ============================================================================
# Tier 1 · 结构性 layout(5 个 · 重写 · 含 TEAM 插画 + LOGO 位 + Speaker 位)
# ============================================================================

def make_cover(prs: _Pres, title: str, subtitle: str, *,
               prepared_by: str = "", date: str = "", version: str = "",
               project_code: str = "", classification: str = "") -> Slide:
    """复刻 native page-1:LOGO 位 + 大字标题 + 橙红圆角副标 + TEAM 插画 + Speaker 位"""
    s = _blank_slide(prs)

    # 左上 LOGO 占位
    _logo_placeholder(s, Inches(0.5), Inches(0.4))

    # 右上 classification 徽标(可选)
    if classification:
        _text_xy(s, Inches(10.5), Inches(0.4), Inches(2.3), Inches(0.4),
                 classification.upper(), size=10, bold=True,
                 color=PRIMARY, align=PP_ALIGN.RIGHT)

    # 主标题(右对齐,大字)
    _text_xy(s, Inches(2.5), Inches(1.3), Inches(10.3), Inches(1.1),
             title, size=42, bold=True, color=H.GRAY_900,
             align=PP_ALIGN.RIGHT)

    # 副标题(橙红圆角框)
    if subtitle:
        sub_w, sub_h = Inches(3.5), Inches(0.55)
        sub_x = Emu(H.SLIDE_W - Inches(0.5) - sub_w)
        sub_y = Inches(2.55)
        frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    sub_x, sub_y, sub_w, sub_h)
        frame.adjustments[0] = 0.4
        frame.fill.background()
        frame.line.color.rgb = PRIMARY
        frame.line.width = Pt(1.5)
        tf = frame.text_frame
        H.fix_textbox_margins(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = subtitle
        H.set_font(r, name=FONT_HEADER, size=14, color=PRIMARY)

    # 底部:TEAM 插画(占左下大块 · 新 aspect 1.46 比 template_a 2.15 更高)
    if TEAM_ILLUSTRATION.exists():
        H.embed_picture(s, str(TEAM_ILLUSTRATION),
                        Inches(0.3), Inches(3.3), height=Inches(3.7))

    # 右下 Speaker name + Designed by(占位文字)
    speaker_lines = []
    if prepared_by or date:
        speaker_lines.append(f"{prepared_by}{' · ' + date if date else ''}")
    if version or project_code:
        meta = " · ".join(v for v in (version, project_code) if v)
        speaker_lines.append(meta)
    if not speaker_lines:
        speaker_lines = ["Speaker name and title", "Designed by 你的团队"]

    for i, line in enumerate(speaker_lines):
        _text_xy(s, Inches(9.0), Emu(Inches(6.4) + i * Inches(0.3)),
                 Inches(3.8), Inches(0.3), line,
                 size=11, color=H.GRAY_500, align=PP_ALIGN.RIGHT)

    return s


def make_toc(prs: _Pres, sections: list[str],
             section_captions: list[str] | None = None,
             caption: str | None = None) -> Slide:
    """复刻 native page-2:双栏,左大字"目录" + 右 N 行菱形序号 + 标题。

    最多显 5 行(模板视觉舒适区上限);若给 ≥ 6 section,只取前 5。
    中间项(n // 2)以橙红菱形高亮"重点章节"。

    Args:
        sections: 章节标题列表
        section_captions: 可选,每章节下方的副标说明。None / 空字符串 → 该章不渲染副标。
                          长度需 ≤ len(sections);超出位用 None 占位。
        caption: 可选,左栏"目录"下方说明。None → 用动态计数 fallback。
    """
    s = _blank_slide(prs)
    n = min(len(sections), 5)

    # === 左栏:CONTENTS 大字 + 描述 ===
    left_box_w = Inches(4.5)
    left_box_h = Inches(1.4)
    left_box_x = Inches(0.55)
    left_box_y = Inches(2.5)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              left_box_x, left_box_y, left_box_w, left_box_h)
    box.fill.background()
    box.line.color.rgb = H.GRAY_300
    box.line.width = Pt(0.5)

    _text_xy(s, Emu(left_box_x + Inches(0.3)), Emu(left_box_y + Inches(0.2)),
             Emu(left_box_w - Inches(0.6)), Inches(1.0),
             "目录", size=44, bold=True, color=H.GRAY_900,
             align=PP_ALIGN.RIGHT)

    caption_text = caption if caption else f"本份内容覆盖 {n} 个核心章节,请见右侧目录。"
    _text_xy(s, left_box_x, Emu(left_box_y + left_box_h + Inches(0.15)),
             left_box_w, Inches(0.4),
             caption_text, size=12,
             color=H.GRAY_500, align=PP_ALIGN.RIGHT)

    # === 右栏:N 行菱形序号 + 标题 + 可选描述 ===
    right_x = Inches(6.0)
    right_w = Inches(7.0)
    section_h = Inches(0.95)
    total_h = Emu(section_h * n + Inches(0.2) * (n - 1))
    start_y = Emu(Inches(7.5 / 2) - total_h // 2)
    captions = list(section_captions or [])

    for i, sec in enumerate(sections[:n]):
        row_y = Emu(start_y + i * (section_h + Inches(0.2)))

        diamond_size = Inches(0.7)
        diamond_x = right_x
        diamond_y = Emu(row_y + Inches(0.1))
        diamond = s.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                      diamond_x, diamond_y,
                                      diamond_size, diamond_size)
        is_highlight = (i == n // 2)
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = PRIMARY if is_highlight else H.GRAY_500
        diamond.line.fill.background()
        tf = diamond.text_frame
        H.fix_textbox_margins(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{i+1:02d}"
        H.set_font(r, name=FONT_NUM, size=14, bold=True, color=H.WHITE)

        title_x = Emu(diamond_x + diamond_size + Inches(0.3))
        title_w = Emu(right_w - diamond_size - Inches(0.3))
        sec_caption = captions[i] if i < len(captions) else None
        title_y_offset = Inches(0.05) if sec_caption else Inches(0.2)
        _text_xy(s, title_x, Emu(row_y + title_y_offset), title_w, Inches(0.35),
                 sec, size=18, bold=True, color=H.GRAY_900)

        if sec_caption:
            _text_xy(s, title_x, Emu(row_y + Inches(0.45)),
                     title_w, Inches(0.3),
                     sec_caption,
                     size=10, color=H.GRAY_500)

        if i < n - 1:
            line_y = Emu(row_y + section_h + Inches(0.05))
            H.rect(s, title_x, line_y, title_w, Pt(0.3), H.GRAY_300)

    return s


def make_section_divider(prs: _Pres, num: int | str, title: str,
                          sub_caption: str = "") -> Slide:
    """复刻 native page-3:左侧橙红方块 + 大数字 /NN + 右侧 TEAM 插画 + 下方浅蓝带 + Section Header

    Args:
        num: 章节序号(int 渲染为 /01;str 直接前置 /)
        title: 章节标题
        sub_caption: 可选,标题下方过渡说明 / 桥句。空字符串 → 不渲染。
    """
    s = _blank_slide(prs)

    # 左侧橙红方块
    block_w, block_h = Inches(2.5), Inches(3.7)
    block_x, block_y = Inches(0.6), Inches(0.4)
    H.rect(s, block_x, block_y, block_w, block_h, PRIMARY)

    # 大数字 /NN
    num_text = f"/{num:02d}" if isinstance(num, int) else f"/{num}"
    _text_xy(s, block_x, Emu(block_y + Inches(2.0)),
             block_w, Inches(1.5), num_text,
             size=80, bold=True, color=H.WHITE,
             font=FONT_NUM, align=PP_ALIGN.CENTER)

    # 右侧 TEAM 插画
    if TEAM_ILLUSTRATION.exists():
        H.embed_picture(s, str(TEAM_ILLUSTRATION),
                        Inches(6.0), Inches(0.5), height=Inches(3.5))

    # 下方浅蓝带
    band_y = Inches(4.5)
    band_h = Inches(0.9)
    H.rect(s, Inches(0), band_y, H.SLIDE_W, band_h, LIGHT_BLUE)

    # Section Header
    _text_xy(s, Inches(0.6), Inches(5.8),
             Inches(8.0), Inches(0.6),
             title, size=32, bold=True, color=H.GRAY_900)

    # 过渡桥句(空字符串则不渲染)
    if sub_caption:
        _text_xy(s, Inches(0.6), Inches(6.5),
                 Inches(8.0), Inches(0.7),
                 sub_caption,
                 size=12, color=H.GRAY_500)

    return s


def make_cards(prs: _Pres, title: str, cards: list[dict[str, Any]]) -> Slide:
    """复刻 native page-15 风格:每卡左上圆形 icon(橙红/深蓝交替),标题下方居中,body 居中"""
    s = _blank_slide(prs)

    # 标题(顶部左对齐)
    _text_xy(s, Inches(0.55), Inches(0.55),
             Inches(12.2), Inches(0.8),
             title, size=32, bold=True, color=PRIMARY_DEEP)
    # 标题下细线
    H.rect(s, Inches(0.55), Inches(1.45),
           Emu(H.SLIDE_W - Inches(1.1)), Pt(0.5), H.GRAY_300)

    # 卡片行
    region = L.content_region()
    row_y = Emu(region.y + Inches(0.5))
    cols = L.columns(L.Box(x=region.x, y=row_y,
                            w=region.w, h=Inches(3.3)), len(cards))

    for i, (col, card) in enumerate(zip(cols, cards)):
        # 圆形 icon(橙红/深蓝交替)
        icon_size = Inches(0.9)
        icon_x = Emu(col.x + col.w // 2 - icon_size // 2)
        icon_y = col.y
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     icon_x, icon_y, icon_size, icon_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else ACCENT
        circle.line.fill.background()
        # icon 中心文字
        tf = circle.text_frame
        H.fix_textbox_margins(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = card.get("title", "")[:1] if card.get("title") else str(i + 1)
        H.set_font(r, name=FONT_NUM, size=24, bold=True, color=H.WHITE)

        # 标题(icon 下方居中)
        title_y = Emu(icon_y + icon_size + Inches(0.3))
        _text_xy(s, col.x, title_y, col.w, Inches(0.5),
                 card.get("title", ""), size=20, bold=True,
                 color=H.GRAY_900, align=PP_ALIGN.CENTER)

        # body
        body_y = Emu(title_y + Inches(0.6))
        _text_xy(s, col.x, body_y, col.w, Inches(1.5),
                 card.get("body", ""), size=14, color=H.GRAY_700,
                 align=PP_ALIGN.CENTER)

    return s


def make_closing(prs: _Pres, subtitle: str = "",
                 next_steps: list[dict[str, str]] | None = None) -> Slide:
    """复刻 native page-32:LOGO + 大字 Thanks + Slogan + Speaker + TEAM 插画右侧 + 橙红渐变带底"""
    s = _blank_slide(prs)

    # 右上 LOGO
    _logo_placeholder(s, Inches(11.8), Inches(0.4))

    # 左侧:大字 Thanks
    _text_xy(s, Inches(0.6), Inches(1.5),
             Inches(7.0), Inches(1.0),
             "Thanks", size=64, bold=True, color=H.GRAY_900)

    # Slogan
    if subtitle:
        _text_xy(s, Inches(0.6), Inches(2.5),
                 Inches(7.0), Inches(0.7),
                 subtitle, size=28, bold=True, color=H.GRAY_900)

    # Speaker name / Next steps
    speaker_lines = []
    if next_steps:
        speaker_lines = [f"{i+1}. {step.get('action', '')}"
                         for i, step in enumerate(next_steps[:3])]
    else:
        speaker_lines = ["Speaker name and title", "Designed by 你的团队"]

    for i, line in enumerate(speaker_lines):
        _text_xy(s, Inches(0.6), Emu(Inches(3.6) + i * Inches(0.4)),
                 Inches(7.0), Inches(0.4), line,
                 size=14, color=H.GRAY_700)

    # 底部橙红带(左半 · 跟 TEAM 插画并排,模拟 native 橙红渐变底带)
    band_y = Inches(4.8)
    band_h = Inches(1.5)
    H.rect(s, Inches(0), band_y, Inches(7.0), band_h, PRIMARY)

    # 右侧 TEAM 插画
    if TEAM_ILLUSTRATION.exists():
        H.embed_picture(s, str(TEAM_ILLUSTRATION),
                        Inches(6.5), Inches(2.8), height=Inches(4.2))

    return s


# ============================================================================
# Tier 2 · 内容 layout(8 个 · 沿用 tech_blue 结构 + 切 template_training 调色板)
# ============================================================================

def make_single_focus(prs: _Pres, *, big_text: str = "", big_number: str = "",
                      explanation: str = "") -> Slide:
    """大数字 120pt 橙红 + 副标 36pt 深蓝 + 解释 18pt 灰 — 模板最强视觉冲击页"""
    s = _blank_slide(prs)
    region = L.content_region()
    blocks = L.stack(region, [Inches(1.6), Inches(0.8), Inches(0.5)],
                     gap=Inches(0.2), align="middle")
    _text_box(s, blocks[0], big_number, size=120, bold=True, color=PRIMARY,
              font=FONT_NUM, align=PP_ALIGN.CENTER)
    _text_box(s, blocks[1], big_text, size=36, bold=True, color=PRIMARY_DEEP,
              align=PP_ALIGN.CENTER)
    _text_box(s, blocks[2], explanation, size=18, color=H.GRAY_700,
              align=PP_ALIGN.CENTER)
    return s


def make_compare(prs: _Pres, title: str, items: list[dict[str, Any]]) -> Slide:
    """N 列对比表 — header bar 风。recommended=True 用橙红 header + 浅橙 body 高亮。"""
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    block_h = Inches(4.6) if H.is_handout() else Inches(3.5)
    row = L.stack(region, [block_h], align="middle")[0]
    cols = L.columns(row, len(items), gap=Inches(0.15))
    header_h = Inches(0.7)
    body_size = 14 if H.is_handout() else 16
    for col, item in zip(cols, items):
        is_recommended = bool(item.get("recommended", False))
        header_fill = PRIMARY if is_recommended else H.GRAY_300
        header_color = H.WHITE if is_recommended else PRIMARY_DEEP
        body_fill = PRIMARY_TINT if is_recommended else H.WHITE
        body_color = PRIMARY_DEEP if is_recommended else H.GRAY_700

        H.rect(s, col.x, col.y, col.w, header_h, header_fill)
        h_tb = s.shapes.add_textbox(col.x, col.y, col.w, header_h)
        H.fix_textbox_margins(h_tb.text_frame)
        h_tb.text_frame.word_wrap = True
        hp = h_tb.text_frame.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run(); hr.text = item["title"]
        H.set_font(hr, name=FONT_HEADER, size=18, bold=True, color=header_color)

        body_y = col.y + header_h
        body_h = block_h - header_h
        body_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col.x, body_y,
                                         col.w, body_h)
        body_shape.fill.solid(); body_shape.fill.fore_color.rgb = body_fill
        body_shape.line.color.rgb = H.GRAY_300
        body_shape.line.width = Pt(0.75)

        if is_recommended:
            # badge 作为 "stamp" 跨 header/body 边界,半上半下,
            # 避免 title 换行溢出 header 时与 body 内部 badge 重叠
            badge_w = Inches(0.5)
            badge_x = col.x + col.w - badge_w - Inches(0.05)
            badge_y = col.y + header_h - badge_w // 2
            badge = s.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y,
                                        badge_w, badge_w)
            badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT
            H.no_line(badge)
            b_tb = s.shapes.add_textbox(badge_x, badge_y, badge_w, badge_w)
            H.fix_textbox_margins(b_tb.text_frame)
            bp = b_tb.text_frame.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run(); br.text = "✓"
            H.set_font(br, name=FONT_HEADER, size=16, bold=True, color=H.WHITE)

        body_box = L.Box(col.x + Inches(0.25), body_y + Inches(0.25),
                          col.w - Inches(0.5), body_h - Inches(0.5))
        _text_box(s, body_box, item["body"], size=body_size, color=body_color)
    return s


def make_compare_pk(prs: _Pres, title: str, left: dict[str, str],
                    right: dict[str, str]) -> Slide:
    """对决式两选一 — 左橙红 vs 右深蓝(ACCENT),中间深蓝 VS 圆。"""
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    vs_diameter = Inches(1.4)
    gap = Inches(0.4)
    side_w = (region.w - vs_diameter - gap * 2) // 2

    block_h = Inches(4.8) if H.is_handout() else Inches(3.8)
    block_y = region.y + (region.h - block_h) // 2

    body_size = 16 if H.is_handout() else 18

    def _side(x: int, side: dict[str, str], is_left: bool) -> None:
        accent_color = PRIMARY if is_left else ACCENT
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, block_y,
                                  side_w, block_h)
        bg.fill.solid(); bg.fill.fore_color.rgb = H.GRAY_50
        bg.line.color.rgb = H.GRAY_300; bg.line.width = Pt(0.75)
        bg.adjustments[0] = 0.03
        H.rect(s, x, block_y, side_w, Inches(0.08), accent_color)
        title_box = L.Box(x + Inches(0.4), block_y + Inches(0.35),
                            side_w - Inches(0.8), Inches(0.9))
        _text_box(s, title_box, side["title"], size=28, bold=True,
                  color=PRIMARY_DEEP, align=PP_ALIGN.CENTER)
        body_box = L.Box(x + Inches(0.4), block_y + Inches(1.35),
                           side_w - Inches(0.8), block_h - Inches(1.6))
        _text_box(s, body_box, side["body"], size=body_size,
                  color=H.GRAY_700, align=PP_ALIGN.CENTER)

    _side(region.x, left, is_left=True)
    _side(region.x + side_w + gap + vs_diameter + gap, right, is_left=False)

    vs_x = region.x + side_w + gap
    vs_y = block_y + (block_h - vs_diameter) // 2
    vs_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, vs_x, vs_y,
                                     vs_diameter, vs_diameter)
    vs_circle.fill.solid(); vs_circle.fill.fore_color.rgb = PRIMARY_DEEP
    H.no_line(vs_circle)
    vs_tb = s.shapes.add_textbox(vs_x, vs_y, vs_diameter, vs_diameter)
    H.fix_textbox_margins(vs_tb.text_frame)
    vp = vs_tb.text_frame.paragraphs[0]
    vp.alignment = PP_ALIGN.CENTER
    vr = vp.add_run(); vr.text = "VS"
    H.set_font(vr, name=FONT_HEADER, size=36, bold=True, color=H.WHITE)
    return s


def make_matrix_2x2(prs: _Pres, title: str, x_axis: dict[str, str],
                    y_axis: dict[str, str],
                    quadrants: list[dict[str, Any]]) -> Slide:
    """BCG 2×2 矩阵 — 高亮象限用浅橙底 + 橙红边框。"""
    s = _blank_slide(prs)
    _add_title(s, title)

    matrix_x = Inches(2.4)
    matrix_y = Inches(1.7)
    matrix_w = Inches(10.0)
    matrix_h = Inches(4.7)
    cell_w = matrix_w // 2
    cell_h = matrix_h // 2

    body_size = 14 if H.is_handout() else 12

    positions = {
        "tl": (matrix_x, matrix_y),
        "tr": (matrix_x + cell_w, matrix_y),
        "bl": (matrix_x, matrix_y + cell_h),
        "br": (matrix_x + cell_w, matrix_y + cell_h),
    }
    for q in quadrants:
        pos = q.get("pos")
        if pos not in positions:
            raise ValueError(f"quadrant.pos 必须是 tl/tr/bl/br,得到 {pos!r}")
        qx, qy = positions[pos]
        highlight = bool(q.get("highlight", False))
        fill = PRIMARY_TINT if highlight else H.WHITE
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, qx, qy, cell_w, cell_h)
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = PRIMARY if highlight else H.GRAY_300
        rect.line.width = Pt(1.5) if highlight else Pt(0.75)
        title_box = L.Box(qx + Inches(0.25), qy + Inches(0.2),
                            cell_w - Inches(0.5), Inches(0.5))
        _text_box(s, title_box, q["title"], size=18, bold=True,
                  color=PRIMARY_DEEP if highlight else H.GRAY_900)
        body_box = L.Box(qx + Inches(0.25), qy + Inches(0.85),
                           cell_w - Inches(0.5), cell_h - Inches(1.0))
        _text_box(s, body_box, q.get("body", ""), size=body_size,
                  color=PRIMARY_DEEP if highlight else H.GRAY_700)

    axis_y = matrix_y + matrix_h + Inches(0.1)
    x_low = L.Box(matrix_x, axis_y, cell_w, Inches(0.35))
    x_high = L.Box(matrix_x + cell_w, axis_y, cell_w, Inches(0.35))
    _text_box(s, x_low, x_axis.get("low", ""), size=12, bold=True,
              color=H.GRAY_700, align=PP_ALIGN.CENTER)
    _text_box(s, x_high, x_axis.get("high", ""), size=12, bold=True,
              color=H.GRAY_700, align=PP_ALIGN.CENTER)

    y_axis_x = Inches(0.55)
    y_axis_w = matrix_x - y_axis_x - Inches(0.15)
    y_high = L.Box(y_axis_x, matrix_y, y_axis_w, Inches(0.4))
    y_low = L.Box(y_axis_x, matrix_y + matrix_h - Inches(0.4), y_axis_w,
                    Inches(0.4))
    _text_box(s, y_high, "↑ " + y_axis.get("high", ""), size=12, bold=True,
              color=H.GRAY_700)
    _text_box(s, y_low, "↓ " + y_axis.get("low", ""), size=12, bold=True,
              color=H.GRAY_700)
    return s


def make_bullet_list(prs: _Pres, title: str, items: list[Any]) -> Slide:
    """要点列表 — 橙红竖条 ▎ 项目符号 + 灰色 body。"""
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    bullet_size = 14 if H.is_handout() else 18
    line_factor = 1.6 if H.is_handout() else 1.45

    if len(items) <= 6:
        bullet_size = max(bullet_size, 22 if H.is_handout() else 26)
        line_factor = 1.8

    if all(isinstance(it, str) for it in items):
        line_h = Emu(int(Pt(bullet_size) * line_factor))
        block = L.stack(region, [Emu(line_h * len(items))], align="middle")[0]
        H.bullets(s, block.x, block.y, block.w, block.h, items=items,
                  size=bullet_size, accent_color=PRIMARY, body_color=H.GRAY_700)
        return s

    # mixed/icon mode
    line_h_emu = Pt(bullet_size * line_factor * 1.5)
    total_h = line_h_emu * len(items)
    block = L.stack(region, [Emu(int(total_h))], align="middle")[0]
    icon_w = Inches(0.45)
    for i, it in enumerate(items):
        y = block.y + Emu(int(line_h_emu * i))
        if isinstance(it, dict):
            text = it.get("text", "")
            icon_char = it.get("icon", "▎")
        else:
            text = str(it); icon_char = "▎"
        icon_str = H.ICONS.get(icon_char, icon_char)
        H.icon(s, block.x, y, bullet_size, icon_str, color=PRIMARY,
               box_size=icon_w)
        text_box = L.Box(block.x + icon_w + Inches(0.1), y,
                          block.w - icon_w - Inches(0.1),
                          Emu(int(line_h_emu)))
        _text_box(s, text_box, text, size=bullet_size, color=H.GRAY_700)
    return s


def make_table(prs: _Pres, title: str, headers: list[str],
               rows: list[list[str]]) -> Slide:
    """表格 — header 深蓝实色 + 斑马纹浅橙。"""
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    H.table_modern(s, region.x, region.y, region.w, region.h,
                   headers=headers, rows=rows,
                   header_fill=PRIMARY_DEEP, header_color=H.WHITE,
                   zebra=PRIMARY_TINT, font_size=14)
    return s


def make_pic_text(prs: _Pres, title: str, image_path: str,
                  points: list[dict[str, str]]) -> Slide:
    """左图右文 — 卡 accent 条用橙红。"""
    s = _blank_slide(prs)
    _add_title(s, title)
    left, right = L.split(L.content_region(), 0.42)
    H.embed_picture(s, image_path, left.x, left.y, box_w=left.w, box_h=left.h)
    rboxes = L.rows(right, len(points))
    body_size = 11 if H.is_handout() else 13
    body_box_h = Inches(0.9) if H.is_handout() else Inches(0.5)
    for rb, p in zip(rboxes, points):
        H.card(s, rb.x, rb.y, rb.w, rb.h, fill=H.WHITE,
               border=H.GRAY_300, accent=PRIMARY)
        inner = L.inset(rb, Inches(0.25), Inches(0.12))
        parts = L.stack(inner, [Inches(0.4), body_box_h], gap=Inches(0.05),
                        align="top")
        _text_box(s, parts[0], p["title"], size=16, bold=True, color=PRIMARY_DEEP)
        _text_box(s, parts[1], p["body"], size=body_size, color=H.GRAY_700)
    return s


def make_summary(prs: _Pres, conclusions: list[str],
                 title: str = "核心结论") -> Slide:
    """总结 — 橙红编号方块(整体竖排居中)+ 右侧结论文字。"""
    s = _blank_slide(prs)
    _add_title(s, title, size=36, color=PRIMARY_DEEP)
    region = L.content_region()
    text_size = 16 if H.is_handout() else 22
    unit_h = Inches(1.2) if H.is_handout() else Inches(0.9)
    gap = Inches(0.25)
    total_h = unit_h * len(conclusions) + gap * (len(conclusions) - 1)
    start_y = region.y + (region.h - total_h) // 2
    num_w = Inches(1.0)
    text_x = region.x + num_w + Inches(0.3)
    text_w = region.w - num_w - Inches(0.3)
    for i, c in enumerate(conclusions):
        y = start_y + (unit_h + gap) * i
        H.rect(s, region.x, y, num_w, unit_h, PRIMARY)
        n_tb = s.shapes.add_textbox(region.x, y, num_w, unit_h)
        H.fix_textbox_margins(n_tb.text_frame)
        n_tb.text_frame.vertical_anchor = 3
        pn = n_tb.text_frame.paragraphs[0]
        pn.alignment = PP_ALIGN.CENTER
        rn = pn.add_run(); rn.text = str(i + 1)
        H.set_font(rn, name=FONT_NUM, size=36, bold=True, color=H.WHITE)
        t_tb = s.shapes.add_textbox(text_x, y, text_w, unit_h)
        H.fix_textbox_margins(t_tb.text_frame)
        t_tb.text_frame.word_wrap = True
        t_tb.text_frame.vertical_anchor = 3
        rt = t_tb.text_frame.paragraphs[0].add_run(); rt.text = c
        H.set_font(rt, name=FONT_HEADER, size=text_size, color=H.GRAY_900)
    return s


# ============================================================================
# Tier 3 · Visual Patterns(library/visual-patterns/ 命中实现)
# ============================================================================

def make_timeline_band_3(prs: _Pres, title: str,
                          segments: list[dict[str, Any]]) -> Slide:
    """visual-pattern: timeline-band-3 · 3 段色块时间轴

    3 个等宽矩形色块横排(橙-灰-橙交替),每块下方居中放时段标签,
    色块上方/下方交错放标题 + 描述。

    Args:
        title: 页标题
        segments: list of 3 dicts {label: "W1", period: "周一-周五",
                                   title: "工程试点", body: "..."}
    """
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    n = min(len(segments), 3)
    if n == 0:
        return s

    band_h = Inches(0.9)
    band_y = region.y + (region.h - band_h) // 2
    gap = Inches(0.2)
    band_w = (region.w - gap * (n - 1)) // n

    body_size = 14 if H.is_handout() else 16
    title_size = 18

    for i in range(n):
        seg = segments[i]
        bx = region.x + (band_w + gap) * i
        color = PRIMARY if i % 2 == 0 else ACCENT
        H.rect(s, bx, band_y, band_w, band_h, color)

        label = seg.get("label", "")
        if label:
            _text_xy(s, bx, band_y, band_w, band_h,
                     label, size=28, bold=True, color=H.WHITE,
                     font=FONT_NUM, align=PP_ALIGN.CENTER)

        period = seg.get("period", "")
        if period:
            _text_xy(s, bx, band_y + band_h + Inches(0.1),
                     band_w, Inches(0.35),
                     period, size=11, color=H.GRAY_500,
                     align=PP_ALIGN.CENTER)

        seg_title = seg.get("title", "")
        seg_body = seg.get("body", "")
        above = (i % 2 == 0)
        if above:
            title_y = band_y - Inches(1.45)
            body_y = band_y - Inches(1.05)
        else:
            title_y = band_y + band_h + Inches(0.6)
            body_y = band_y + band_h + Inches(1.0)

        if seg_title:
            _text_xy(s, bx, title_y, band_w, Inches(0.4),
                     seg_title, size=title_size, bold=True,
                     color=PRIMARY_DEEP, align=PP_ALIGN.CENTER)
        if seg_body:
            body_box = L.Box(bx + Inches(0.15), body_y,
                              band_w - Inches(0.3), Inches(1.0))
            _text_box(s, body_box, seg_body, size=body_size,
                      color=H.GRAY_700, align=PP_ALIGN.CENTER)

    return s


def make_tri_pyramid_4sub_3(prs: _Pres, title: str,
                              items: list[dict[str, Any]]) -> Slide:
    """visual-pattern: tri-pyramid-4sub-3 · 大三角拆 4 子三角 + 3 编号说明

    大等边三角形(浅蓝边框 + 蓝填充),内部切成 4 个子三角(中间倒置浅色 "洞"),
    左下 / 右下 / 顶角放 1./2./3. 编号 + 标题 + 描述。

    Args:
        title: 页标题
        items: list of 3 dicts {title: "工程", body: "..."}
              顺序:items[0] = 左下 / items[1] = 右下 / items[2] = 顶
    """
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    n = min(len(items), 3)
    if n == 0:
        return s

    # 大三角占左 ~45%,右侧放 3 个编号说明
    tri_zone_w = Emu(region.w * 0.42)
    tri_zone_x = region.x
    tri_zone_y = region.y
    tri_zone_h = region.h

    # 等边三角形 — 居中
    tri_h = min(tri_zone_h - Inches(0.4), Emu(tri_zone_w * 0.866))  # √3/2
    tri_w = Emu(tri_h / 0.866)
    tri_x = tri_zone_x + (tri_zone_w - tri_w) // 2
    tri_y = tri_zone_y + (tri_zone_h - tri_h) // 2

    # 大三角(浅蓝边 + 主色填充)
    big_tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  tri_x, tri_y, tri_w, tri_h)
    big_tri.fill.solid()
    big_tri.fill.fore_color.rgb = PRIMARY_DEEP
    big_tri.line.color.rgb = LIGHT_BLUE
    big_tri.line.width = Pt(2.0)

    # 中央倒置子三角("洞")— 边长为大三角一半,居中
    sub_w = tri_w // 2
    sub_h = tri_h // 2
    inv_x = tri_x + (tri_w - sub_w) // 2
    inv_y = tri_y + tri_h - sub_h - sub_h // 4
    inv_tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  inv_x, inv_y, sub_w, sub_h)
    inv_tri.fill.solid()
    inv_tri.fill.fore_color.rgb = LIGHT_BLUE
    H.no_line(inv_tri)
    inv_tri.rotation = 180

    # 右侧 3 个编号说明
    right_x = tri_zone_x + tri_zone_w + Inches(0.4)
    right_w = region.x + region.w - right_x
    right_gap = Inches(0.25)
    item_h = (region.h - right_gap * (n - 1)) // n
    body_size = 13 if H.is_handout() else 14

    for i in range(n):
        item = items[i]
        iy = region.y + (item_h + right_gap) * i

        # 编号方块
        num_size = Inches(0.7)
        H.rect(s, right_x, iy, num_size, num_size, PRIMARY)
        n_tb = s.shapes.add_textbox(right_x, iy, num_size, num_size)
        H.fix_textbox_margins(n_tb.text_frame)
        n_tb.text_frame.vertical_anchor = 3
        pn = n_tb.text_frame.paragraphs[0]
        pn.alignment = PP_ALIGN.CENTER
        rn = pn.add_run(); rn.text = f"{i + 1:02d}"
        H.set_font(rn, name=FONT_NUM, size=16, bold=True, color=H.WHITE)

        # 标题
        title_x = right_x + num_size + Inches(0.2)
        title_w = right_w - num_size - Inches(0.2)
        _text_xy(s, title_x, iy, title_w, Inches(0.4),
                 item.get("title", ""), size=18, bold=True,
                 color=PRIMARY_DEEP)

        # 描述
        body_box = L.Box(title_x, iy + Inches(0.45),
                          title_w, item_h - Inches(0.45))
        _text_box(s, body_box, item.get("body", ""),
                  size=body_size, color=H.GRAY_700)

    return s


def make_cards_flag_3(prs: _Pres, title: str,
                       cards: list[dict[str, Any]]) -> Slide:
    """visual-pattern: cards-flag-3 · 3 张旗帜风卡(顶部撕角 + 圆 icon)

    3 张矩形卡(浅蓝 / 浅橙 / 浅绿渐变背景),卡顶右角"撕角"折叠效果;
    卡顶中央放彩色圆 + (可选) icon 字符;圆下方放粗体标题 + 短下划线,再下方灰色描述。

    Args:
        title: 页标题
        cards: list of 3 dicts {title, body, icon?: "★" / "1" / emoji}
    """
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    n = min(len(cards), 3)
    if n == 0:
        return s

    # 三色 tint 调色(模板亲缘 · 偏浅)
    BG_TINTS = [
        RGBColor(0xE8, 0xF1, 0xFB),  # 浅蓝
        RGBColor(0xFD, 0xE8, 0xDF),  # 浅橙
        RGBColor(0xE6, 0xF3, 0xEC),  # 浅绿
    ]
    ICON_COLORS = [ACCENT, PRIMARY, RGBColor(0x40, 0x96, 0x94)]  # 深蓝 / 橙红 / 青绿

    gap = Inches(0.25)
    card_h = Inches(4.6) if H.is_handout() else Inches(3.6)
    card_w = (region.w - gap * (n - 1)) // n
    cards_y = region.y + (region.h - card_h) // 2
    body_size = 13 if H.is_handout() else 14
    tear = Inches(0.35)  # 撕角尺寸

    for i in range(n):
        card = cards[i]
        cx = region.x + (card_w + gap) * i
        bg = BG_TINTS[i % 3]
        icon_color = ICON_COLORS[i % 3]

        # 主卡形 — 五边形 (撕角效果) 用 freeform
        from pptx.util import Emu as _Emu
        from pptx.enum.shapes import MSO_CONNECTOR
        # 用 rect 加一个右上角小三角"撕"
        H.rect(s, cx, cards_y, card_w, card_h, bg)
        # 右上角撕角(白色三角覆盖)
        tear_tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                       cx + card_w - tear, cards_y,
                                       tear, tear)
        tear_tri.fill.solid()
        tear_tri.fill.fore_color.rgb = H.WHITE
        H.no_line(tear_tri)
        tear_tri.rotation = 180  # 翻转露出右上"缺角"
        # 撕角折痕线(深一点的边)
        crease = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                     cx + card_w - tear, cards_y,
                                     tear, tear)
        crease.fill.background()
        crease.line.color.rgb = H.GRAY_300
        crease.line.width = Pt(0.5)
        crease.rotation = 180

        # 顶部圆 icon
        icon_d = Inches(0.85)
        icon_x = cx + (card_w - icon_d) // 2
        icon_y = cards_y + Inches(0.4)
        icon_circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                          icon_x, icon_y, icon_d, icon_d)
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = icon_color
        H.no_line(icon_circle)
        icon_char = card.get("icon", f"{i + 1:02d}")
        ic_tb = s.shapes.add_textbox(icon_x, icon_y, icon_d, icon_d)
        H.fix_textbox_margins(ic_tb.text_frame)
        ic_tb.text_frame.vertical_anchor = 3
        ic_p = ic_tb.text_frame.paragraphs[0]
        ic_p.alignment = PP_ALIGN.CENTER
        ic_r = ic_p.add_run(); ic_r.text = icon_char
        H.set_font(ic_r, name=FONT_NUM, size=22, bold=True, color=H.WHITE)

        # 标题
        title_y = icon_y + icon_d + Inches(0.2)
        _text_xy(s, cx + Inches(0.2), title_y,
                 card_w - Inches(0.4), Inches(0.4),
                 card.get("title", ""), size=17, bold=True,
                 color=PRIMARY_DEEP, align=PP_ALIGN.CENTER)

        # 下划线
        ul_w = Inches(0.6)
        ul_x = cx + (card_w - ul_w) // 2
        ul_y = title_y + Inches(0.45)
        H.rect(s, ul_x, ul_y, ul_w, Pt(2.5), icon_color)

        # 描述
        body_y = ul_y + Inches(0.25)
        body_h = cards_y + card_h - body_y - Inches(0.3)
        body_box = L.Box(cx + Inches(0.3), body_y,
                          card_w - Inches(0.6), body_h)
        _text_box(s, body_box, card.get("body", ""),
                  size=body_size, color=H.GRAY_700,
                  align=PP_ALIGN.CENTER)

    return s
