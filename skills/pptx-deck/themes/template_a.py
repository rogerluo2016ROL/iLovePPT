"""template_a 主题 —— iSlide 商业模板的部分复刻(Phase 2,~70% 视觉相似度)

# 复刻原则
- Tier 1(重写,占视觉权重 80%): make_cover / make_toc / make_section_divider / make_closing
- Tier 2(覆盖局部): make_cards 改成圆形 icon
- Tier 3(继承 tech_blue 同名实现): single_focus / compare / bullet_list / table / pic_text / summary

# 模板 visual identity
- 双主色:橙红 #EF5938 + 深蓝 #0B2A4A
- 标志性元素:TEAM 四人扁平插画(_assets/template_template_a/team_illustration.png)
- LOGO 占位 / Speaker name 占位 / Designed by 署名
- section_divider 渐变方块 + 大数字 + 浅蓝带

# 复刻限制(不到 100% 的部分)
- SmartArt 类(辐射图 / 4 象限装饰)= 0% 复刻(iLovePPT 11 layout 不含)
- iSlide footer 水印 = 不复制(模板作者签名,非视觉规范)
- 复杂渐变 / 多层装饰 = 简化复刻
"""
import sys
from pathlib import Path
from typing import Any

# Fallback for direct import outside pytest
_helpers_path = str(Path(__file__).parent.parent.parent / "pptx")
if _helpers_path not in sys.path:
    sys.path.insert(0, _helpers_path)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as _Pres
from pptx.slide import Slide
from pptx.util import Emu, Inches, Length, Pt

import helpers as H
import layout as L

# ===== 模板视觉常量 =====
FONT_HEADER = H.FONT_CN
FONT_BODY   = H.FONT_CN
FONT_NUM    = H.FONT_NUM

# 双主色
PRIMARY      = RGBColor(0xEF, 0x59, 0x38)   # 橙红(模板 accent1)
PRIMARY_DEEP = RGBColor(0x0B, 0x2A, 0x4A)   # 深蓝
PRIMARY_TINT = RGBColor(0xFD, 0xE5, 0xDE)   # 浅橙
LIGHT_BLUE   = RGBColor(0xC8, 0xE4, 0xF8)   # 浅蓝带(section_divider 用)
ACCENT       = RGBColor(0x08, 0x59, 0x86)   # 深蓝 accent(模板 accent2)

# 资产路径(指仓库根的 _assets/template_template_a/)
_REPO_ROOT = Path(__file__).resolve().parent.parent.parent.parent
TEAM_ILLUSTRATION = _REPO_ROOT / "_assets" / "template_template_a" / "team_illustration.png"


def _blank_slide(prs: _Pres) -> Slide:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide: Slide, x: Length, y: Length, w: Length, h: Length,
          text: str, *, size: int, bold: bool = False, color=None,
          align=PP_ALIGN.LEFT, font=None) -> None:
    if color is None:
        color = H.GRAY_900
    if font is None:
        font = FONT_HEADER
    tb = slide.shapes.add_textbox(x, y, w, h)
    H.fix_textbox_margins(tb.text_frame)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    H.set_font(r, name=font, size=size, bold=bold, color=color)


def _logo_placeholder(slide: Slide, x: Length, y: Length) -> None:
    """LOGO 占位框 — 浅灰描边小框 + 'LOGO' 文字"""
    w, h = Inches(0.9), Inches(0.4)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.background()
    rect.line.color.rgb = H.GRAY_300
    rect.line.width = Pt(0.5)
    tf = rect.text_frame
    H.fix_textbox_margins(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "LOGO"
    H.set_font(r, name=FONT_HEADER, size=10, color=H.GRAY_500)


# ============================================================================
# Tier 1 · 重写 4 大 layout
# ============================================================================

def make_cover(prs: _Pres, title: str, subtitle: str, *,
               prepared_by: str = "", date: str = "", version: str = "",
               project_code: str = "", classification: str = "") -> Slide:
    """复刻 native page-1:LOGO 位 + 大字标题 + 橙红圆角副标 + TEAM 插画 + Speaker 位"""
    s = _blank_slide(prs)

    # 左上 LOGO 占位
    _logo_placeholder(s, Inches(0.5), Inches(0.4))

    # 右上 classification 徽标(可选)
    if classification:
        _text(s, Inches(10.5), Inches(0.4), Inches(2.3), Inches(0.4),
              classification.upper(), size=10, bold=True,
              color=PRIMARY, align=PP_ALIGN.RIGHT)

    # 主标题(右对齐,大字)
    _text(s, Inches(2.5), Inches(1.3), Inches(10.3), Inches(1.1),
          title, size=42, bold=True, color=H.GRAY_900,
          align=PP_ALIGN.RIGHT)

    # 副标题(橙红圆角框)
    if subtitle:
        sub_w, sub_h = Inches(3.5), Inches(0.55)
        sub_x = Emu(H.SLIDE_W - Inches(0.5) - sub_w)
        sub_y = Inches(2.55)
        frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    sub_x, sub_y, sub_w, sub_h)
        frame.adjustments[0] = 0.4   # 大圆角
        frame.fill.background()
        frame.line.color.rgb = PRIMARY
        frame.line.width = Pt(1.5)
        tf = frame.text_frame
        H.fix_textbox_margins(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = subtitle
        H.set_font(r, name=FONT_HEADER, size=14, color=PRIMARY)

    # 底部:TEAM 插画(占左下大块)
    if TEAM_ILLUSTRATION.exists():
        H.embed_picture(s, str(TEAM_ILLUSTRATION),
                        Inches(0.3), Inches(3.5), height=Inches(3.5))

    # 右下 Speaker name + Designed by(占位文字)
    speaker_lines = []
    if prepared_by or date:
        speaker_lines.append(f"{prepared_by}{' · ' + date if date else ''}")
    if version or project_code:
        meta = " · ".join(v for v in (version, project_code) if v)
        speaker_lines.append(meta)
    if not speaker_lines:
        speaker_lines = ["Speaker name and title", "Designed by 你的团队"]

    for i, line in enumerate(speaker_lines):
        _text(s, Inches(9.0), Emu(Inches(6.4) + i * Inches(0.3)),
              Inches(3.8), Inches(0.3), line,
              size=11, color=H.GRAY_500, align=PP_ALIGN.RIGHT)

    return s


def make_toc(prs: _Pres, sections: list[str]) -> Slide:
    """复刻 native page-2:双栏,左大字 CONTENTS + 描述;右侧 N 行带菱形序号 + 标题 + 描述"""
    s = _blank_slide(prs)

    # === 左栏:CONTENTS 大字 + 描述 ===
    left_box_w = Inches(4.5)
    left_box_h = Inches(1.4)
    left_box_x = Inches(0.55)
    left_box_y = Inches(2.5)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              left_box_x, left_box_y, left_box_w, left_box_h)
    box.fill.background()
    box.line.color.rgb = H.GRAY_300
    box.line.width = Pt(0.5)

    _text(s, Emu(left_box_x + Inches(0.3)), Emu(left_box_y + Inches(0.2)),
          Emu(left_box_w - Inches(0.6)), Inches(1.0),
          "目录", size=44, bold=True, color=H.GRAY_900,
          align=PP_ALIGN.RIGHT)

    _text(s, left_box_x, Emu(left_box_y + left_box_h + Inches(0.15)),
          left_box_w, Inches(0.4),
          "本份内容覆盖 N 个核心章节,请见右侧目录。", size=12,
          color=H.GRAY_500, align=PP_ALIGN.RIGHT)

    # === 右栏:N 行菱形序号 + 标题 + 描述 ===
    right_x = Inches(6.0)
    right_w = Inches(7.0)
    n = min(len(sections), 5)
    section_h = Inches(0.95)
    total_h = Emu(section_h * n + Inches(0.2) * (n - 1))
    start_y = Emu(Inches(7.5 / 2) - total_h // 2)

    for i, sec in enumerate(sections[:n]):
        row_y = Emu(start_y + i * (section_h + Inches(0.2)))

        diamond_size = Inches(0.7)
        diamond_x = right_x
        diamond_y = Emu(row_y + Inches(0.1))
        diamond = s.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                      diamond_x, diamond_y,
                                      diamond_size, diamond_size)
        is_highlight = (i == n // 2)
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = PRIMARY if is_highlight else H.GRAY_500
        diamond.line.fill.background()
        tf = diamond.text_frame
        H.fix_textbox_margins(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{i+1:02d}"
        H.set_font(r, name=FONT_NUM, size=14, bold=True, color=H.WHITE)

        title_x = Emu(diamond_x + diamond_size + Inches(0.3))
        title_w = Emu(right_w - diamond_size - Inches(0.3))
        _text(s, title_x, row_y, title_w, Inches(0.35),
              sec, size=18, bold=True, color=H.GRAY_900)

        _text(s, title_x, Emu(row_y + Inches(0.4)),
              title_w, Inches(0.3),
              "Copy paste fonts. Choose the only option to retain text.",
              size=10, color=H.GRAY_500)

        if i < n - 1:
            line_y = Emu(row_y + section_h + Inches(0.05))
            H.rect(s, title_x, line_y, title_w, Pt(0.3), H.GRAY_300)

    return s


def make_section_divider(prs: _Pres, num: int | str, title: str) -> Slide:
    """复刻 native page-3:左侧橙红方块 + 大数字 /NN + 右侧 TEAM 插画 + 下方浅蓝带 + Section Header"""
    s = _blank_slide(prs)

    # 左侧橙红方块
    block_w, block_h = Inches(2.5), Inches(3.7)
    block_x, block_y = Inches(0.6), Inches(0.4)
    H.rect(s, block_x, block_y, block_w, block_h, PRIMARY)

    # 大数字 /NN
    num_text = f"/{num:02d}" if isinstance(num, int) else f"/{num}"
    _text(s, block_x, Emu(block_y + Inches(2.0)),
          block_w, Inches(1.5), num_text,
          size=80, bold=True, color=H.WHITE,
          font=FONT_NUM, align=PP_ALIGN.CENTER)

    # 右侧 TEAM 插画
    if TEAM_ILLUSTRATION.exists():
        H.embed_picture(s, str(TEAM_ILLUSTRATION),
                        Inches(6.5), Inches(0.7), height=Inches(3.3))

    # 下方浅蓝带
    band_y = Inches(4.5)
    band_h = Inches(0.9)
    H.rect(s, Inches(0), band_y, H.SLIDE_W, band_h, LIGHT_BLUE)

    # Section Header
    _text(s, Inches(0.6), Inches(5.8),
          Inches(8.0), Inches(0.6),
          title, size=32, bold=True, color=H.GRAY_900)

    # Supporting text
    _text(s, Inches(0.6), Inches(6.5),
          Inches(8.0), Inches(0.4),
          "Supporting text. When you copy & paste, choose \"keep text only\" option.",
          size=12, color=H.GRAY_500)

    return s


def make_closing(prs: _Pres, subtitle: str = "",
                 next_steps: list[dict[str, str]] | None = None) -> Slide:
    """复刻 native page-32:LOGO + 大字 Thanks + Slogan + Speaker + TEAM 插画右侧 + 橙红渐变带底"""
    s = _blank_slide(prs)

    # 右上 LOGO
    _logo_placeholder(s, Inches(11.8), Inches(0.4))

    # 左侧:大字 Thanks
    _text(s, Inches(0.6), Inches(1.5),
          Inches(7.0), Inches(1.0),
          "Thanks", size=64, bold=True, color=H.GRAY_900)

    # Slogan
    if subtitle:
        _text(s, Inches(0.6), Inches(2.5),
              Inches(7.0), Inches(0.7),
              subtitle, size=28, bold=True, color=H.GRAY_900)

    # Speaker name / Next steps
    speaker_lines = []
    if next_steps:
        speaker_lines = [f"{i+1}. {step.get('action', '')}"
                         for i, step in enumerate(next_steps[:3])]
    else:
        speaker_lines = ["Speaker name and title", "Designed by 你的团队"]

    for i, line in enumerate(speaker_lines):
        _text(s, Inches(0.6), Emu(Inches(3.6) + i * Inches(0.4)),
              Inches(7.0), Inches(0.4), line,
              size=14, color=H.GRAY_700)

    # 底部橙红带
    band_y = Inches(4.8)
    band_h = Inches(1.5)
    H.rect(s, Inches(0), band_y, Inches(7.0), band_h, PRIMARY)

    # 右侧 TEAM 插画
    if TEAM_ILLUSTRATION.exists():
        H.embed_picture(s, str(TEAM_ILLUSTRATION),
                        Inches(6.5), Inches(2.8), height=Inches(4.2))

    return s


# ============================================================================
# Tier 2 · 改 make_cards 加圆形 icon
# ============================================================================

def make_cards(prs: _Pres, title: str, cards: list[dict[str, str]]) -> Slide:
    """覆盖 tech_blue 实现:卡片左侧改成圆形 icon(native page-15 风格)"""
    s = _blank_slide(prs)

    # 标题(顶部左对齐)
    _text(s, Inches(0.55), Inches(0.55),
          Inches(12.2), Inches(0.8),
          title, size=32, bold=True, color=H.GRAY_900)
    # 标题下细线
    H.rect(s, Inches(0.55), Inches(1.45),
           Emu(H.SLIDE_W - Inches(1.1)), Pt(0.5), H.GRAY_300)

    # 卡片行
    region = L.content_region()
    row_y = Emu(region.y + Inches(0.5))
    cols = L.columns(L.Box(x=region.x, y=row_y,
                            w=region.w, h=Inches(3.3)), len(cards))

    for i, (col, card) in enumerate(zip(cols, cards)):
        # 圆形 icon(橙红/深蓝交替)
        icon_size = Inches(0.9)
        icon_x = Emu(col.x + col.w // 2 - icon_size // 2)
        icon_y = col.y
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     icon_x, icon_y, icon_size, icon_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else ACCENT
        circle.line.fill.background()
        # icon 中心文字
        tf = circle.text_frame
        H.fix_textbox_margins(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = card.get("title", "")[:1] if card.get("title") else str(i + 1)
        H.set_font(r, name=FONT_NUM, size=24, bold=True, color=H.WHITE)

        # 标题(icon 下方居中)
        title_y = Emu(icon_y + icon_size + Inches(0.3))
        _text(s, col.x, title_y, col.w, Inches(0.5),
              card.get("title", ""), size=20, bold=True,
              color=H.GRAY_900, align=PP_ALIGN.CENTER)

        # body
        body_y = Emu(title_y + Inches(0.6))
        _text(s, col.x, body_y, col.w, Inches(1.5),
              card.get("body", ""), size=14, color=H.GRAY_700,
              align=PP_ALIGN.CENTER)

    return s


# ============================================================================
# Tier 3 · 其他 layout 从 tech_blue 继承(不改)
# ============================================================================

from themes.tech_blue import (
    make_single_focus,
    make_compare,
    make_bullet_list,
    make_table,
    make_pic_text,
    make_summary,
)
