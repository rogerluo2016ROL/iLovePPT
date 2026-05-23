"""iLovePPT build.py —— 机械构建器：deck_plan.json → .pptx + PNG。

用法：python3 build.py deck_plan.json [--no-render]

deck_plan.json schema：{theme, output, slides: [{layout, ...fields}, ...]}
智能部分（brief→deck_plan、视觉自检）由 Claude 按文档流程做,不在本文件。
"""
import sys
import json
import shutil
import subprocess
from pathlib import Path
from types import ModuleType
from typing import Any

from pptx import Presentation

HERE = Path(__file__).parent
for _p in [str(HERE.parent / "pptx"), str(HERE)]:
    if _p not in sys.path:
        sys.path.insert(0, _p)

import helpers as H
from themes import tech_blue as _tech_blue

THEMES: dict[str, ModuleType] = {"tech_blue": _tech_blue}

# 需要页脚 + 页码的 layout(规范:visual-qa.md §页脚 / 页码完整性)。
# cover / section_divider / closing 不计入页码。
FOOTERED_LAYOUTS: frozenset[str] = frozenset({
    "toc", "single_focus", "compare", "cards",
    "bullet_list", "table", "pic_text", "summary",
})


# ----- load_plan -----

def load_plan(path: str | Path) -> dict[str, Any]:
    """读 + 校验 deck_plan.json。记录 _plan_dir 供相对 output 解析。"""
    p = Path(path)
    data = json.loads(p.read_text(encoding="utf-8"))
    for field in ("theme", "output", "slides"):
        if field not in data:
            raise ValueError(f"deck_plan 缺字段: {field}")
    if not isinstance(data["slides"], list) or not data["slides"]:
        raise ValueError("deck_plan.slides 必须是非空 list")
    for i, slide in enumerate(data["slides"], 1):
        if "layout" not in slide:
            raise ValueError(f"deck_plan 第 {i} 页缺 layout 字段")
    data["_plan_dir"] = str(p.resolve().parent)
    return data


# ----- load_theme -----

def _extract_design_tokens(pptx_path: str) -> dict[str, Any]:
    """从 .pptx 提取主色(accent1)与中文字体(master ea typeface)。best-effort。"""
    from lxml import etree
    tokens: dict[str, Any] = {}
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return tokens
    try:
        from pptx.oxml.ns import qn
        if prs.slide_masters:
            done = False
            for ph in prs.slide_masters[0].placeholders:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(qn("a:rPr"))
                        if rPr is not None:
                            ea = rPr.find(qn("a:ea"))
                            if ea is not None and ea.get("typeface"):
                                tokens["font_header"] = ea.get("typeface")
                                tokens["font_body"] = ea.get("typeface")
                                done = True
                                break
                    if done:
                        break
                if done:
                    break
    except Exception:
        pass
    try:
        from pptx.dml.color import RGBColor
        for part in prs.part.package.iter_parts():
            pn = part.partname
            if "theme" in pn and pn.endswith(".xml"):
                root = etree.fromstring(part.blob)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                accent1 = root.find(".//a:accent1//a:srgbClr", ns)
                if accent1 is not None:
                    hx = accent1.get("val", "")
                    if len(hx) == 6:
                        tokens["primary"] = RGBColor(
                            int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                        break
    except Exception:
        pass
    return tokens


def _extract_theme_from_pptx(pptx_path: str) -> ModuleType:
    """从用户 .pptx 提取主色与字体,派生临时主题模块。

    从 tech_blue.py 源码加载全新模块实例,再用提取的 token 覆盖
    FONT_* / PRIMARY 属性。token 提取 best-effort,未提取到的保留默认。
    """
    import importlib.util
    tokens = _extract_design_tokens(pptx_path)
    import themes.tech_blue as base_module
    base_path = Path(base_module.__file__)
    out_name = f"extracted_{Path(pptx_path).stem}"
    spec = importlib.util.spec_from_file_location(out_name, base_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法从 {base_path} 加载主题")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    if "font_header" in tokens:
        mod.FONT_HEADER = tokens["font_header"]
        mod.FONT_BODY = tokens.get("font_body", tokens["font_header"])
    if "primary" in tokens:
        mod.PRIMARY = tokens["primary"]
    font_status = tokens.get("font_header", "默认 Microsoft YaHei")
    color_status = tokens.get("primary", "默认 tech_blue 主色")
    print(f"  从模板提取主题: {out_name}")
    print(f"     字体: {font_status}")
    print(f"     主色: {color_status}")
    return mod


def _repo_templates_dir() -> Path:
    """iLovePPT 仓库根的 templates/ 目录(`<repo>/skills/pptx-deck/build.py`
    → `<repo>/templates/`)。.resolve() 处理符号链接场景。"""
    return Path(__file__).resolve().parent.parent.parent / "templates"


def _find_template(name: str, plan_dir: str | None = None) -> Path | None:
    """按短名查找 .pptx 模板。

    优先级:
      1. <plan_dir>/templates/<name>.pptx  (deck 项目专属)
      2. <repo>/templates/<name>.pptx      (全局共享)

    找到返回 Path,找不到返回 None。
    """
    candidates: list[Path] = []
    if plan_dir:
        candidates.append(Path(plan_dir) / "templates" / f"{name}.pptx")
    candidates.append(_repo_templates_dir() / f"{name}.pptx")
    for p in candidates:
        if p.exists():
            return p
    return None


def _list_available_templates() -> list[str]:
    """返回 <repo>/templates/ 下所有 .pptx 短名(不含扩展名)。"""
    tdir = _repo_templates_dir()
    if not tdir.exists():
        return []
    return sorted(p.stem for p in tdir.glob("*.pptx"))


def load_theme(theme_id: str, plan_dir: str | None = None) -> ModuleType:
    """解析 theme_id 到 theme 模块。

    Args:
        theme_id: 三种形式之一
            - 内置 theme 名(如 "tech_blue")
            - 短名(如 "company_a")—— 查找 templates/<name>.pptx
            - 路径(含 "/" 或以 ".pptx" 结尾)—— 直接当 .pptx 路径
        plan_dir: deck plan 所在目录,影响相对路径解析 + 短名查找优先级
    """
    if theme_id in THEMES:
        return THEMES[theme_id]
    # 含 / 或以 .pptx 结尾 → 当路径处理
    if str(theme_id).endswith(".pptx") or "/" in str(theme_id):
        path = Path(theme_id).expanduser()
        if not path.is_absolute() and plan_dir:
            path = (Path(plan_dir) / path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"theme .pptx 不存在: {path}")
        return _extract_theme_from_pptx(str(path))
    # 短名 → 查 templates/
    found = _find_template(theme_id, plan_dir)
    if found is not None:
        return _extract_theme_from_pptx(str(found))
    # 未找到 → 列可用的帮用户排错
    available = _list_available_templates()
    available_str = ", ".join(available) if available else "(空,把 .pptx 放进 templates/)"
    raise ValueError(
        f"未知 theme: {theme_id!r}. "
        f"内置: tech_blue. "
        f"templates/ 可用: {available_str}. "
        f"或直接给 .pptx 绝对/相对路径。"
    )


# ----- build_deck -----

def build_deck(plan: dict[str, Any]) -> Path:
    """按 deck_plan 逐 slide 调 make_*,存 .pptx,返回输出路径。

    自动处理 3 个 cross-cutting 字段(build.py 集中负责,theme 不感知):
    - **footer**: 内容页(FOOTERED_LAYOUTS)统一加分隔线 + "N / TOTAL" + 可选元数据
    - **footer_meta**(plan 顶层): classification / project / version,显示在 footer 左侧
    - **source**(slide 级): 数据 slide 的引文,渲染在 footer 上方
    """
    theme = load_theme(plan["theme"], plan.get("_plan_dir"))
    prs = Presentation()
    prs.slide_width = H.SLIDE_W
    prs.slide_height = H.SLIDE_H

    footer_meta = plan.get("footer_meta", {}) or {}

    # 预扫:算 footer 页总数(用于 "N / TOTAL" 的 TOTAL)
    total_footered = sum(
        1 for s in plan["slides"] if s["layout"] in FOOTERED_LAYOUTS
    )
    footer_idx = 0

    for i, slide in enumerate(plan["slides"], 1):
        layout = slide["layout"]
        fn = getattr(theme, f"make_{layout}", None)
        if fn is None:
            raise ValueError(f"第 {i} 页未知 layout: {layout}（theme 无 make_{layout}）")
        # 弹出 cross-cutting 字段,不传给 make_* fn(避免 TypeError)
        fields = {k: v for k, v in slide.items() if k != "layout"}
        source = fields.pop("source", None)
        try:
            fn(prs, **fields)
        except TypeError as e:
            raise ValueError(f"第 {i} 页 layout={layout}: {e}") from e

        # 数据引文(source)→ footer 上方
        if source:
            H.source_citation(prs.slides[-1], source)

        # 页脚 + 页码 + footer_meta
        if layout in FOOTERED_LAYOUTS and total_footered > 0:
            footer_idx += 1
            H.footer(
                prs.slides[-1], footer_idx, total_footered,
                classification=footer_meta.get("classification"),
                project=footer_meta.get("project"),
                version=footer_meta.get("version"),
            )

    out = Path(plan["output"]).expanduser()
    if not out.is_absolute():
        out = (Path(plan["_plan_dir"]) / out).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


# ----- render -----

def render(pptx_path: str | Path, out_dir: str | Path) -> list[Path]:
    """soffice → PDF → pdftoppm → 逐页 PNG。返回 PNG 路径列表。"""
    if shutil.which("soffice") is None:
        raise RuntimeError("soffice 未安装。请: brew install --cask libreoffice")
    if shutil.which("pdftoppm") is None:
        raise RuntimeError("pdftoppm 未安装。请: brew install poppler")
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             str(pptx_path), "--outdir", str(out_dir)],
            check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"soffice 转 PDF 失败: {e.stderr}") from e
    pdf = out_dir / (pptx_path.stem + ".pdf")
    if not pdf.exists():
        raise RuntimeError(f"soffice 跑了但未产 PDF: {pdf}")
    try:
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "120", str(pdf), str(out_dir / "page")],
            check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"pdftoppm 转 PNG 失败: {e.stderr}") from e
    return sorted(out_dir.glob("page-*.jpg"))


# ----- CLI -----

def main(argv: list[str]) -> None:
    if not argv:
        sys.exit("用法: python3 build.py deck_plan.json [--no-render]")
    plan_path = argv[0]
    do_render = "--no-render" not in argv
    plan = load_plan(plan_path)
    out = build_deck(plan)
    print(f"已生成 {out}")
    if do_render:
        render_dir = out.parent / (out.stem + "_render")
        pngs = render(out, render_dir)
        print(f"已渲染 {len(pngs)} 页 → {render_dir}")


if __name__ == "__main__":
    main(sys.argv[1:])
