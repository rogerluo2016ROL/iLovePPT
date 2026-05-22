"""iLovePPT pptx-deck — 端到端 workflow 骨架。

用法：
    python3 workflow.py path/to/brief.yaml

流程：parse_brief → load_theme → plan_diagrams → outline → per-slide(generate + render + vision_check) → save
vision_check 当前实现为：导出 PNG,打印路径,默认接受。
真实使用时由 Claude 调本脚本后逐张看图,出 issue JSON 再 fix。
plan_diagrams 当前为关键词匹配骨架；真实判断由 Claude 按 diagram-planning.md 做。
"""
import sys, subprocess, tempfile
from pathlib import Path
from types import ModuleType
from typing import Any

import yaml
from pptx import Presentation
from pptx.presentation import Presentation as _Pres
from pptx.slide import Slide

HERE = Path(__file__).parent
# Fallback for direct script execution (pytest uses pyproject.toml pythonpath)
for _p in [str(HERE.parent / "pptx"), str(HERE)]:
    if _p not in sys.path:
        sys.path.insert(0, _p)

import helpers as H
from themes import tech_blue as T
import re


# ----- 1. parse_brief -----

REQUIRED: set[str] = {"title", "outline", "theme", "output"}


def parse_brief(path: str | Path) -> dict[str, Any]:
    with open(path) as f:
        data = yaml.safe_load(f)
    missing = REQUIRED - set(data)
    if missing:
        raise ValueError(f"brief 缺字段: {missing}")
    data.setdefault("subtitle", "")
    data.setdefault("page_count_target", None)
    data.setdefault("key_points", [])
    data.setdefault("reference_pptx", None)
    return data


# ----- 2. load_theme -----

THEMES: dict[str, ModuleType] = {"tech_blue": T}


def _extract_design_tokens(pptx_path: str) -> dict[str, Any]:
    """从 .pptx 提取 design token (主色 + 字体)。提取失败时返回空 dict,由调用方回退默认值。"""
    from pptx import Presentation as _Pres2
    from lxml import etree

    tokens: dict[str, Any] = {}
    try:
        prs = _Pres2(pptx_path)
    except Exception:
        return tokens

    # 1. 字体 — 取 slide master 中第一个含 ea typeface 的 run
    try:
        from pptx.oxml.ns import qn
        if prs.slide_masters:
            master = prs.slide_masters[0]
            outer_break = False
            for ph in master.placeholders:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(qn("a:rPr"))
                        if rPr is not None:
                            ea = rPr.find(qn("a:ea"))
                            if ea is not None and ea.get("typeface"):
                                tokens["font_header"] = ea.get("typeface")
                                tokens["font_body"] = ea.get("typeface")
                                outer_break = True
                                break
                    if outer_break:
                        break
                if outer_break:
                    break
    except Exception:
        pass

    # 2. 主色 — 从 theme*.xml 读 accent1 srgbClr
    try:
        from pptx.dml.color import RGBColor
        for part in prs.part.package.iter_parts():
            pn = part.partname
            if "theme" in pn and pn.endswith(".xml"):
                root = etree.fromstring(part.blob)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                accent1 = root.find(".//a:accent1//a:srgbClr", ns)
                if accent1 is not None:
                    hex_val = accent1.get("val", "")
                    if len(hex_val) == 6:
                        tokens["primary"] = RGBColor(
                            int(hex_val[0:2], 16),
                            int(hex_val[2:4], 16),
                            int(hex_val[4:6], 16),
                        )
                        break
    except Exception:
        pass

    return tokens


def _ingest_template(pptx_path: str) -> ModuleType:
    """从用户 .pptx 学风格,派生一个临时主题模块。

    实现：从 tech_blue.py 源码加载一个全新模块实例（与缓存的 themes.tech_blue
    隔离）,再用提取的 token 覆盖其 FONT_* / PRIMARY 属性。make_* 函数在调用时
    按所在模块的 __dict__ 解析全局名,故 import 后直接设属性即可生效——无需
    字符串替换源码（那种做法在 tech_blue 改为 SSOT 别名后已失效）。

    token 提取为 best-effort：未提取到的字段保留 tech_blue 默认值。
    """
    import importlib.util

    tokens = _extract_design_tokens(pptx_path)

    import themes.tech_blue as base_module
    base_path = Path(base_module.__file__)
    out_name = f"ingested_{Path(pptx_path).stem}"

    spec = importlib.util.spec_from_file_location(out_name, base_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法从 {base_path} 加载 ingested theme")
    ingested = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(ingested)

    # 用提取的 token 覆盖主题常量（设属性即可,见 docstring）
    if "font_header" in tokens:
        ingested.FONT_HEADER = tokens["font_header"]
        ingested.FONT_BODY = tokens.get("font_body", tokens["font_header"])
    if "primary" in tokens:
        ingested.PRIMARY = tokens["primary"]

    print(f"  ingested theme: {out_name}")
    print(f"     fonts: {tokens.get('font_header', '(default)')}")
    print(f"     primary: {tokens.get('primary', '(default)')}")
    return ingested


def load_theme(theme_id: str) -> ModuleType:
    if theme_id in THEMES:
        return THEMES[theme_id]
    if str(theme_id).endswith(".pptx"):
        if not Path(theme_id).exists():
            raise FileNotFoundError(f"theme .pptx 文件不存在: {theme_id}")
        return _ingest_template(theme_id)
    raise ValueError(f"未知 theme: {theme_id}")


# ----- 3. plan_diagrams -----

# 触发配图的关键词信号 → (图类型, 推荐工具)。判断细则见 diagram-planning.md。
_DIAGRAM_SIGNALS: list[tuple[str, str, str]] = [
    (r"架构|系统组成|分层|模块|组件|微服务|技术栈|拓扑", "arch_diagram", "draw.io"),
    (r"流程|步骤|阶段|时序|先后|环节|pipeline|workflow", "flow", "mermaid"),
    (r"趋势|增长|占比|对比|数据|指标|百分比|\d+\s*%", "chart", "matplotlib"),
    (r"关系|依赖|交互|连接|映射|角色边界", "simple_relation", "pptx-native"),
]


def plan_diagrams(brief: dict[str, Any]) -> list[dict[str, Any]]:
    """图层规划：读完 brief 后,判断哪些章节适合配图。

    骨架版用关键词匹配（_DIAGRAM_SIGNALS）扫描每个 outline 章节标题。
    真实运行时 Claude 按 diagram-planning.md 的决策规则做语义判断,
    覆盖本函数的输出（关键词命中只是粗筛提示,不是最终决定）。

    每个 outline 章节最多产出 1 个图建议（命中的第一个信号）。

    Returns: list of {section_idx, section, diagram_type, tool, intent}
    """
    plan: list[dict[str, Any]] = []
    for sec_idx, sec in enumerate(brief["outline"], 1):
        for pattern, dtype, tool in _DIAGRAM_SIGNALS:
            if re.search(pattern, sec):
                plan.append({
                    "section_idx": sec_idx,
                    "section": sec,
                    "diagram_type": dtype,
                    "tool": tool,
                    "intent": f"为「{sec}」配 {dtype}",
                })
                break  # 每章节最多 1 张图
    return plan


# ----- 4. outline → page_specs -----

def estimate_page_count(brief: dict[str, Any]) -> int:
    if brief["page_count_target"]:
        return brief["page_count_target"]
    return int(len(brief["outline"]) * 1.5) + 4


def generate_outline(
    brief: dict[str, Any],
    diagram_plan: list[dict[str, Any]] | None = None,
) -> list[dict[str, Any]]:
    """根据 brief 生成 page_spec list。LLM 在真实运行时会替换此函数。
    本骨架返回固定的简版 outline 跑通 pipeline。

    diagram_plan（plan_diagrams 的输出）中命中的章节,其内容页会带上
    `visual_element` 元数据 {type, tool, intent}。真实运行时 Claude 据此
    调 [[diagram]] skill 出图,并把该页 layout 换成 pic_text。
    """
    by_section = {d["section_idx"]: d for d in (diagram_plan or [])}
    specs = []
    specs.append({"layout": "cover", "title": brief["title"],
                  "subtitle": brief.get("subtitle", "")})
    specs.append({"layout": "toc", "sections": brief["outline"]})
    for i, sec in enumerate(brief["outline"], 1):
        specs.append({"layout": "section_divider", "num": i, "title": sec})
        kp = brief.get("key_points") or [f"{sec} 要点 1", f"{sec} 要点 2"]
        content: dict[str, Any] = {"layout": "bullet_list", "title": sec, "items": kp[:5]}
        if i in by_section:
            d = by_section[i]
            content["visual_element"] = {
                "type": d["diagram_type"],
                "tool": d["tool"],
                "intent": d["intent"],
            }
        specs.append(content)
    specs.append({"layout": "summary",
                  "conclusions": brief.get("key_points",
                                            ["结论 1", "结论 2", "结论 3"])})
    specs.append({"layout": "closing"})  # subtitle 留空,只显大字"谢谢"
    return specs


# ----- 5. generate_slide -----

# 非渲染参数（规划元数据）,generate_slide 调 make_* 前需剔除
_NON_RENDER_KEYS = {"layout", "visual_element"}


def generate_slide(prs: _Pres, spec: dict[str, Any], theme: ModuleType) -> Slide:
    fn = getattr(theme, f"make_{spec['layout']}")
    kwargs = {k: v for k, v in spec.items() if k not in _NON_RENDER_KEYS}
    return fn(prs, **kwargs)


# ----- 6. render_one_slide -----

def render_one_slide(prs: _Pres, idx: int, out_png: str | Path) -> None:
    """导出全 deck PDF,然后 pdftoppm 截第 idx 页。"""
    import shutil  # 局部 import 防止顶部污染
    if shutil.which("soffice") is None:
        raise RuntimeError("soffice 未安装。请: brew install --cask libreoffice")
    if shutil.which("pdftoppm") is None:
        raise RuntimeError("pdftoppm 未安装。请: brew install poppler")
    tmpdir = Path(tempfile.gettempdir()) / "iloveppt_render"
    tmpdir.mkdir(exist_ok=True)
    pptx_tmp = tmpdir / "current.pptx"
    prs.save(str(pptx_tmp))
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             str(pptx_tmp), "--outdir", str(tmpdir)],
            check=True, capture_output=True, text=True,
        )
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"soffice PDF 转换失败: {e.stderr}") from e
    pdf = tmpdir / "current.pdf"
    if not pdf.exists():
        raise RuntimeError(f"soffice 跑了但未产 PDF: {pdf}")
    try:
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "120",
             "-f", str(idx), "-l", str(idx),
             str(pdf), str(tmpdir / "slide_only")],
            check=True, capture_output=True, text=True,
        )
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"pdftoppm 转换失败: {e.stderr}") from e
    candidates = list(tmpdir.glob(f"slide_only-{idx}*.jpg"))
    if not candidates:
        raise RuntimeError(f"渲染未产 jpg,期望第 {idx} 页输出")
    Path(out_png).parent.mkdir(parents=True, exist_ok=True)
    candidates[0].rename(out_png)


def vision_check(image_path: str | Path, intent: str) -> list[dict[str, Any]]:
    """占位：默认接受。真实运行时 Claude 用 Read tool 看图后输出 issue JSON。"""
    print(f"  [vision_check] {image_path} (intent: {intent})")
    return []


def _fix_fontsize_too_large(slide: Slide, issue: dict[str, Any]) -> str:
    """对 slide 中所有 textbox 的 font.size 减 20%（最小 8pt）。"""
    from pptx.util import Pt
    fixed = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is None:
                    continue
                new_size = int(run.font.size.pt * 0.8)
                if new_size < 8:
                    new_size = 8
                run.font.size = Pt(new_size)
                fixed += 1
    return f"font 缩小 20% (影响 {fixed} 个 run)"


def _fix_margin_not_zeroed(slide: Slide, issue: dict[str, Any]) -> str:
    """所有 textbox margin 归零。"""
    from pptx.util import Emu
    fixed = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        fixed += 1
    return f"textbox margin 归零 (影响 {fixed} 个 textbox)"


def _fix_no_action(slide: Slide, issue: dict[str, Any]) -> str:
    """无可机械应用的修复,记录 suggested_fix 供人工处理。"""
    return f"无机械修复策略: {issue.get('suggested_fix', '(空)')}"


def fix_slide(slide: Slide, issues: list[dict[str, Any]]) -> Slide:
    """根据 issues 修 slide。

    支持的机械修复：
    - suggested_fix 含 "字号" / "font" + ("大"|"large"|"small") → 全 slide font 缩小 20%
    - suggested_fix 含 "margin" / "归零" → 全 slide textbox margin 归零
    - 其他关键字 → 打印 suggested_fix 但不修改（限制：视觉 QA 输出为自由文本,
      只能处理有限的关键字模式；复杂修复需人工介入）

    Returns the (possibly modified) slide.
    """
    print(f"  [fix_slide] 应用 {len(issues)} 个修复")
    for issue in issues:
        sf = issue.get("suggested_fix", "").lower()
        # 字号过大/过小：含 "字号" 或 "font" + size-related 词
        if "字号" in sf or ("font" in sf and re.search(r"大|large|过大|small|过小", sf)):
            action = _fix_fontsize_too_large(slide, issue)
        # margin 未归零
        elif "margin" in sf or "归零" in sf:
            action = _fix_margin_not_zeroed(slide, issue)
        else:
            action = _fix_no_action(slide, issue)
        print(f"    - {issue.get('issue', '(unknown)')}: {action}")
    return slide


# ----- 7. main loop -----

def run(brief_path: str | Path) -> tuple[Path, list[dict[str, Any]]]:
    brief = parse_brief(brief_path)
    theme = load_theme(brief["theme"])
    diagram_plan = plan_diagrams(brief)
    if diagram_plan:
        print(f"diagram plan: {len(diagram_plan)} 个章节建议配图")
        for d in diagram_plan:
            print(f"  - 章节 {d['section_idx']}「{d['section']}」"
                  f" → {d['diagram_type']} ({d['tool']})")
    outline = generate_outline(brief, diagram_plan)

    prs = Presentation()
    prs.slide_width = H.SLIDE_W   # SSOT: helpers.py
    prs.slide_height = H.SLIDE_H

    review_needed = []

    for idx, spec in enumerate(outline, 1):
        print(f"slide {idx}/{len(outline)}: {spec['layout']}")
        slide = generate_slide(prs, spec, theme)
        png_path = str(Path(tempfile.gettempdir()) / "iloveppt_render" / f"page_{idx:02d}.png")
        try:
            render_one_slide(prs, idx, png_path)
        except RuntimeError as e:
            print(f"  warning: render failed: {e}; marking review-needed")
            review_needed.append({"idx": idx, "reason": "render_failed"})
            continue

        attempts = 0
        issues = vision_check(png_path, intent=spec["layout"])
        while issues and attempts < 3:
            slide = fix_slide(slide, issues)
            render_one_slide(prs, idx, png_path)
            issues = vision_check(png_path, intent=spec["layout"])
            attempts += 1
        if issues:
            review_needed.append({"idx": idx, "reason": "vision_unresolved",
                                  "issues": issues})

    # output 相对路径按 brief 文件所在目录解析（而非 CWD）,
    # 这样从任意目录跑 workflow.py 产物落点都一致。
    out = Path(brief["output"]).expanduser()
    if not out.is_absolute():
        out = (Path(brief_path).resolve().parent / out).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"\nDone: {out}")
    if review_needed:
        print(f"Warning: {len(review_needed)} pages need review:")
        for r in review_needed:
            print(f"  - page {r['idx']}: {r['reason']}")
    return out, review_needed


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit("用法: python3 workflow.py path/to/brief.yaml")
    run(sys.argv[1])
