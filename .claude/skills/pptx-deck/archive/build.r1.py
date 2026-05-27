"""iLovePPT build.py —— 机械构建器：deck_plan.json → .pptx + PNG。

用法：python3 build.py deck_plan.json [--no-render]

deck_plan.json schema：{theme, output, slides: [{layout, ...fields}, ...]}
智能部分（brief→deck_plan、视觉自检）由 Claude 按文档流程做,不在本文件。
"""
import copy
import sys
import json
import shutil
import subprocess
from pathlib import Path
from types import ModuleType
from typing import Any

from pptx import Presentation
from pptx.shapes.group import GroupShape

HERE = Path(__file__).parent
for _p in [str(HERE.parent / "pptx"), str(HERE)]:
    if _p not in sys.path:
        sys.path.insert(0, _p)

import helpers as H
from themes import tech_blue as _tech_blue
from themes import template_golden as _template_golden
from themes import template_training as _template_training

THEMES: dict[str, ModuleType] = {
    "tech_blue": _tech_blue,
    "template_training": _template_training,
}

# Themes with Python layout implementations but **not** registered in THEMES.
# These are ingested .pptx templates (library/pptx-templates/items/<name>/) whose
# theme module supplies layout shapes; color/font tokens come from the .pptx at
# runtime via _extract_theme_from_pptx. Keying off the .pptx stem lets
# `_extract_theme_from_pptx` pick the right base module instead of always
# hardcoding tech_blue.
PPTX_BASE_THEMES: dict[str, ModuleType] = {
    "template_golden": _template_golden,
}

# 需要页脚 + 页码的 layout(规范:visual-qa.md §页脚 / 页码完整性)。
# cover / section_divider / closing 不计入页码。
FOOTERED_LAYOUTS: frozenset[str] = frozenset({
    "toc", "single_focus", "compare", "compare_pk", "matrix_2x2", "cards",
    "bullet_list", "table", "pic_text", "summary",
    "timeline_band_3", "tri_pyramid_4sub_3", "cards_flag_3",
})


# ----- load_plan -----

def load_plan(path: str | Path) -> dict[str, Any]:
    """读 + 校验 deck_plan.json。记录 _plan_dir 供相对 output 解析。"""
    p = Path(path)
    data = json.loads(p.read_text(encoding="utf-8"))
    for field in ("theme", "output", "slides"):
        if field not in data:
            raise ValueError(f"deck_plan 缺字段: {field}")
    if not isinstance(data["slides"], list) or not data["slides"]:
        raise ValueError("deck_plan.slides 必须是非空 list")
    for i, slide in enumerate(data["slides"], 1):
        if "layout" not in slide:
            raise ValueError(f"deck_plan 第 {i} 页缺 layout 字段")
    data["_plan_dir"] = str(p.resolve().parent)
    return data


# ----- load_theme -----

def _extract_design_tokens(pptx_path: str) -> dict[str, Any]:
    """从 .pptx 提取扩展 design token(accent1-6 + 字号阶梯).

    返回字段(全 best-effort,缺则不在 dict 里):
    - font_header / font_body: master ea typeface
    - primary: RGBColor(accent1)
    - accent_2/3/4/5/6: RGBColor(accent2-6)
    - dk1 / lt1: RGBColor(主文本色 / 主背景色)
    - title_size_pt / body_size_pt: master 默认字号(int pt)
    """
    from lxml import etree
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    tokens: dict[str, Any] = {}

    def _hex2rgb(hx: str) -> RGBColor | None:
        if len(hx) == 6:
            return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        return None

    try:
        prs = Presentation(pptx_path)
    except Exception:
        return tokens

    # === master ea typeface + 字号 ===
    try:
        if prs.slide_masters:
            done = False
            for ph in prs.slide_masters[0].placeholders:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(qn("a:rPr"))
                        if rPr is not None:
                            ea = rPr.find(qn("a:ea"))
                            if ea is not None and ea.get("typeface"):
                                tokens["font_header"] = ea.get("typeface")
                                tokens["font_body"] = ea.get("typeface")
                                done = True
                                break
                    if done:
                        break
                if done:
                    break
    except Exception:
        pass

    # === 从 master XML 抽字号(直接读 slideMaster1.xml) ===
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
              "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        for part in prs.part.package.iter_parts():
            pn = part.partname
            if "slideMaster" in pn and pn.endswith(".xml"):
                root = etree.fromstring(part.blob)
                title_def = root.find(".//p:titleStyle//a:lvl1pPr/a:defRPr", ns)
                if title_def is not None:
                    sz = title_def.get("sz")
                    if sz and sz.isdigit():
                        tokens["title_size_pt"] = int(sz) // 100
                body_def = root.find(".//p:bodyStyle//a:lvl1pPr/a:defRPr", ns)
                if body_def is not None:
                    sz = body_def.get("sz")
                    if sz and sz.isdigit():
                        tokens["body_size_pt"] = int(sz) // 100
                break
    except Exception:
        pass

    # === theme1.xml: accent1-6 + dk1/lt1 ===
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for part in prs.part.package.iter_parts():
            pn = part.partname
            if "theme" in pn and pn.endswith(".xml"):
                root = etree.fromstring(part.blob)
                scheme = root.find(".//a:clrScheme", ns)
                if scheme is None:
                    continue
                for tag, key in [("accent1", "primary"),
                                  ("accent2", "accent_2"),
                                  ("accent3", "accent_3"),
                                  ("accent4", "accent_4"),
                                  ("accent5", "accent_5"),
                                  ("accent6", "accent_6"),
                                  ("dk1", "dk1"), ("lt1", "lt1")]:
                    node = scheme.find(f"a:{tag}", ns)
                    if node is None:
                        continue
                    srgb = node.find(".//a:srgbClr", ns)
                    if srgb is not None:
                        rgb = _hex2rgb(srgb.get("val", ""))
                        if rgb:
                            tokens[key] = rgb
                            continue
                    # dk1/lt1 可能是 sysClr,取 lastClr
                    if tag in ("dk1", "lt1"):
                        sys_clr = node.find("a:sysClr", ns)
                        if sys_clr is not None:
                            rgb = _hex2rgb(sys_clr.get("lastClr", ""))
                            if rgb:
                                tokens[key] = rgb
                break
    except Exception:
        pass

    return tokens


def _extract_theme_from_pptx(pptx_path: str) -> ModuleType:
    """从用户 .pptx 提取主色与字体,派生临时主题模块。

    Base 选择:用 .pptx 文件 stem 查 PPTX_BASE_THEMES(如 template_golden 用其
    专属 module 提供独有 layout),未注册则回落到 tech_blue。再从 base 源码加载
    全新 module 实例,用提取的 token 覆盖 FONT_* / PRIMARY 等属性。
    token 提取 best-effort,未提取到的保留 base 默认。
    """
    import importlib.util
    tokens = _extract_design_tokens(pptx_path)
    stem = Path(pptx_path).stem
    base_module = PPTX_BASE_THEMES.get(stem, _tech_blue)
    base_path = Path(base_module.__file__)
    out_name = f"extracted_{stem}"
    spec = importlib.util.spec_from_file_location(out_name, base_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法从 {base_path} 加载主题")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    # 字体
    if "font_header" in tokens:
        mod.FONT_HEADER = tokens["font_header"]
        mod.FONT_BODY = tokens.get("font_body", tokens["font_header"])
    # 主色 + 次级色(accent2-6 + dk1/lt1)
    if "primary" in tokens:
        mod.PRIMARY = tokens["primary"]
    for key in ("accent_2", "accent_3", "accent_4", "accent_5", "accent_6",
                "dk1", "lt1"):
        if key in tokens:
            setattr(mod, key.upper(), tokens[key])
    # 字号阶梯(若 master 显式定义)
    if "title_size_pt" in tokens:
        mod.TITLE_SIZE_PT = tokens["title_size_pt"]
    if "body_size_pt" in tokens:
        mod.BODY_SIZE_PT = tokens["body_size_pt"]

    font_status = tokens.get("font_header", "默认 Microsoft YaHei")
    color_status = tokens.get("primary", "默认 tech_blue 主色")
    print(f"  从模板提取主题: {out_name}")
    print(f"     字体: {font_status}")
    print(f"     主色: {color_status}")
    extra_accents = [k for k in ("accent_2", "accent_3", "accent_4", "accent_5",
                                  "accent_6") if k in tokens]
    if extra_accents:
        print(f"     次级色: {', '.join(extra_accents)}")
    if "title_size_pt" in tokens or "body_size_pt" in tokens:
        print(f"     字号阶梯: title={tokens.get('title_size_pt', '-')}pt /"
              f" body={tokens.get('body_size_pt', '-')}pt")
    return mod


def _repo_templates_dir() -> Path:
    """仓库根的 library/pptx-templates/_source/ 目录(模板 .pptx 全局共享)。
    build.py 位于 <repo>/.claude/skills/pptx-deck/build.py。"""
    repo_root = Path(__file__).resolve().parents[3]  # build.py → pptx-deck → skills → .claude → repo
    return repo_root / "library" / "pptx-templates" / "_source"


def _find_template(name: str, plan_dir: str | None = None) -> Path | None:
    """按短名查找 .pptx 模板。

    优先级:
      1. <plan_dir>/templates/<name>.pptx                       (deck 项目专属, 向后兼容)
      2. <repo>/library/pptx-templates/_source/<name>.pptx      (全局共享, 新位置)

    找到返回 Path,找不到返回 None。
    """
    candidates: list[Path] = []
    if plan_dir:
        candidates.append(Path(plan_dir) / "templates" / f"{name}.pptx")
    candidates.append(_repo_templates_dir() / f"{name}.pptx")
    for p in candidates:
        if p.exists():
            return p
    return None


def _list_available_templates() -> list[str]:
    """返回 <repo>/library/pptx-templates/_source/ 下所有 .pptx 短名(不含扩展名)。"""
    tdir = _repo_templates_dir()
    if not tdir.exists():
        return []
    return sorted(p.stem for p in tdir.glob("*.pptx"))


def load_theme(theme_id: str, plan_dir: str | None = None) -> ModuleType:
    """解析 theme_id 到 theme 模块。

    Args:
        theme_id: 三种形式之一
            - 内置 theme 名(如 "tech_blue")
            - 短名(如 "company_a")—— 查找 library/pptx-templates/_source/<name>.pptx(全局)
              或 <plan_dir>/templates/<name>.pptx(deck 项目本地)
            - 路径(含 "/" 或以 ".pptx" 结尾)—— 直接当 .pptx 路径
        plan_dir: deck plan 所在目录,影响相对路径解析 + 短名查找优先级
    """
    if theme_id in THEMES:
        return THEMES[theme_id]
    # 含 / 或以 .pptx 结尾 → 当路径处理
    if str(theme_id).endswith(".pptx") or "/" in str(theme_id):
        path = Path(theme_id).expanduser()
        if not path.is_absolute() and plan_dir:
            path = (Path(plan_dir) / path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"theme .pptx 不存在: {path}")
        return _extract_theme_from_pptx(str(path))
    # 短名 → 查 library/pptx-templates/_source/ + <plan_dir>/templates/
    found = _find_template(theme_id, plan_dir)
    if found is not None:
        return _extract_theme_from_pptx(str(found))
    # 未找到 → 列可用的帮用户排错
    available = _list_available_templates()
    available_str = ", ".join(available) if available else "(空,把 .pptx 放进 library/pptx-templates/_source/)"
    raise ValueError(
        f"未知 theme: {theme_id!r}. "
        f"内置: tech_blue. "
        f"library/pptx-templates/_source/ 可用: {available_str}. "
        f"或直接给 .pptx 绝对/相对路径。"
    )


# ----- tier1: template slide reuse -----

# Placeholder text patterns iSlide / template_golden 用作占位符的文字。
# tier1 path 会识别这些 text 并替换成 deck_plan.text_map 提供的内容。
# 注:patterns 按"包含匹配",会命中如 "01.Text here" / "ISLIDE® POWERPOINT"
# / "Speaker name and title" 等扩展形式。
_PLACEHOLDER_PATTERNS = (
    "…text", "...text",                       # 金字塔 tier / 一般占位
    "Text here", "Text Here",                 # 通用("01.Text here" / "02.Text Here" 也命中)
    "Copy paste fonts",                       # 段落 body 占位
    "Supporting text here", "supporting text",  # process_flow 步骤 body
    "www.islide.cc",                          # 页脚网址(主线程会替换 deck footer)
    "ISLIDE", "iSlide",                       # logo 文字残留
    "SUBTITLE HERE", "Subtitle here",         # cover 副标
    "Speaker name", "speaker name",            # 作者位
    "PRESENTATION", "Presentation",            # cover 主标残留
    "LOGO HERE",                              # logo placeholder
    "TITLE", "Title",                          # 通用 title(注意可能误命中,主线程在 title 走 placeholder name match 优先)
)
# Text patterns that are pure placeholder labels (no semantic content) — these
# trigger replacement even when no template-specific text match found
_PLACEHOLDER_BARE_TEXT = {"Text", "TEXT", "01", "02", "03", "04", "05", "06"}


def _walk_shapes(shapes):
    """Recursively yield shapes, descending into groups."""
    for shape in shapes:
        yield shape
        if isinstance(shape, GroupShape):
            yield from _walk_shapes(shape.shapes)


def _is_placeholder_text(text: str) -> bool:
    text = text.strip()
    return any(p in text for p in _PLACEHOLDER_PATTERNS)


def _collect_placeholder_shapes(slide) -> list:
    """Return list of text-bearing placeholder shapes in geometric order.

    Order strategy:
    - First by category (pyramid tiers / cards / sidebar) based on shape size+position
    - Within each category, sorted top-to-bottom then left-to-right
    """
    items: list[tuple] = []
    for shape in _walk_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if not _is_placeholder_text(text):
            continue
        items.append((shape.top or 0, shape.left or 0, shape))
    # Sort top-to-bottom, left-to-right
    items.sort(key=lambda t: (t[0], t[1]))
    return [item[2] for item in items]


def _replace_shape_text(shape, new_text: str, font_ea: str = "Microsoft YaHei",
                         text_color_hex: str | None = None,
                         font_size_pt: int | None = None):
    """Replace shape's text with new_text, preserving first run's font formatting.

    text_color_hex: 可选 "#RRGGBB",覆盖原字色(用于浅底色 tier 强制 dark text)。
    font_size_pt:   可选 int pt,强制覆盖字号(用于 cover 标题中文比英文短需缩小)。
    """
    from lxml import etree
    from pptx.oxml.ns import qn
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_run = first_para.runs[0]
        first_run.text = new_text
        # Ensure ea typeface set for Chinese rendering
        rPr = first_run._r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(first_run._r, qn("a:rPr"))
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font_ea)
        # Override color if specified
        if text_color_hex:
            hex_clean = text_color_hex.lstrip("#")
            # Remove existing solidFill / schemeClr children inside rPr
            for child_name in ("a:solidFill", "a:schemeClr"):
                for child in rPr.findall(qn(child_name)):
                    rPr.remove(child)
            solid_fill = etree.SubElement(rPr, qn("a:solidFill"))
            srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
            srgb.set("val", hex_clean.upper())
        # Override font size if specified
        if font_size_pt:
            rPr.set("sz", str(int(font_size_pt) * 100))
        # Blank other runs in first para
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        # No runs — set text directly
        first_para.text = new_text
    # Blank other paragraphs
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _remove_shape(slide, shape):
    """Remove a shape from a slide (handles top-level shapes only)."""
    sp = shape._element
    parent = sp.getparent()
    if parent is not None:
        parent.remove(sp)


def _copy_slide_from_source(source_slide, target_prs, blank_layout_idx: int = 4):
    """Deep-copy a slide from source_slide into target_prs (appended).

    Uses the source slide's layout (by name match in target_prs.slide_layouts) so
    placeholder size/position/fill inheritance is preserved. Falls back to
    blank_layout_idx if no name match.

    Returns the new slide.
    """
    src_layout_name = source_slide.slide_layout.name
    matched_layout = None
    for lay in target_prs.slide_layouts:
        if lay.name == src_layout_name:
            matched_layout = lay
            break
    if matched_layout is None:
        # Fallback: use specified blank layout
        if blank_layout_idx < len(target_prs.slide_layouts):
            matched_layout = target_prs.slide_layouts[blank_layout_idx]
        else:
            matched_layout = target_prs.slide_layouts[-1]
    new_slide = target_prs.slides.add_slide(matched_layout)
    for shp in list(new_slide.shapes):
        _remove_shape(new_slide, shp)
    for shape in source_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def _walk_tree_paths(shapes, prefix: str = ""):
    """Yield (tree_path, shape) for every shape — top-level + recursively into groups.

    tree_path is a dot-separated index string like "3.0.1" meaning:
      - shape at top-level index 3 (a group)
      - its child at index 0 (a group)
      - its child at index 1 (a leaf or group)

    This matches the path scheme used by placeholder_map.yaml `tree_path` slots.
    Stable across saves because iSlide GroupShape children keep insertion order.
    """
    for i, shape in enumerate(shapes):
        path = f"{prefix}{i}"
        yield path, shape
        if isinstance(shape, GroupShape):
            yield from _walk_tree_paths(shape.shapes, prefix=path + ".")


def _find_shape_by_tree_path(slide, tree_path: str):
    """Locate a shape by tree_path (e.g. "3.0.1.0") via depth-first index navigation."""
    parts = tree_path.split(".")
    shapes = list(slide.shapes)
    cur = None
    for i, p in enumerate(parts):
        try:
            idx = int(p)
        except ValueError:
            return None
        if idx < 0 or idx >= len(shapes):
            return None
        cur = shapes[idx]
        if i < len(parts) - 1:
            if not isinstance(cur, GroupShape):
                return None
            shapes = list(cur.shapes)
    return cur


def _load_placeholder_map(theme: ModuleType, template_page_index: int,
                            plan_dir: str | None) -> dict | None:
    """Load library/pptx-templates/items/<template>/pages/<page-dir>/placeholder_map.yaml.

    Locates page-dir by scanning the pages/ folder for any subdir whose
    placeholder_map.yaml has matching template_page_index. Returns dict or None.
    """
    name = theme.__name__
    if not name.startswith("extracted_"):
        return None
    template_stem = name[len("extracted_"):]
    # Find repo root via build.py location: <repo>/.claude/skills/pptx-deck/build.py
    repo_root = Path(__file__).resolve().parents[3]
    items_dir = repo_root / "library" / "pptx-templates" / "items" / template_stem / "pages"
    if not items_dir.exists():
        return None
    try:
        import yaml as _yaml
    except ImportError:
        return None
    for page_dir in items_dir.iterdir():
        map_path = page_dir / "placeholder_map.yaml"
        if not map_path.exists():
            continue
        try:
            data = _yaml.safe_load(map_path.read_text(encoding="utf-8"))
        except Exception:
            continue
        if isinstance(data, dict) and data.get("template_page_index") == template_page_index:
            return data
    return None


def _apply_tier1_text_map(slide, title: str | None, text_map: dict[str, str],
                            placeholder_map: dict | None = None):
    """On a tier1 slide:
    - replace 标题 placeholder with title (if shape named 标题 1 / Title 1 exists)
    - replace text shapes per placeholder_map (map-driven) OR auto-detect (fallback)
    - remove footer/页号 iSlide placeholders LAST (deck adds its own via H.footer)

    Order matters: text-map replacement happens BEFORE footer removal so tree_path
    indices in placeholder_map stay aligned with the original template's shape order.
    Removing footers first would shift all subsequent top-level indices.
    """
    # Step 1: title rewrite (by name match, doesn't shift indices)
    for shape in slide.shapes:
        name = shape.name or ""
        if shape.has_text_frame and ("标题" in name or "Title" in name):
            if title:
                _replace_shape_text(shape, title)
            break  # Only first title

    # Step 2: text map replacement (BEFORE footer removal — indices must stay original)
    if placeholder_map and placeholder_map.get("slots"):
        _apply_text_map_by_slots(slide, text_map, placeholder_map)
    else:
        _apply_text_map_by_geometry(slide, text_map)

    # Step 3: remove iSlide footer/page-number placeholders (deck adds its own footer)
    shapes_to_remove = []
    for shape in slide.shapes:
        name = shape.name or ""
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "页脚" in name or "Footer" in name or "灯片编号" in name or "Slide Number" in name:
                shapes_to_remove.append(shape)
                continue
            if text in ("www.islide.cc",):
                shapes_to_remove.append(shape)
                continue
    for shp in shapes_to_remove:
        _remove_shape(slide, shp)


def _apply_text_map_by_slots(slide, text_map: dict[str, str], placeholder_map: dict):
    """Map-driven replacement: each slot has tree_path + optional text_color.

    Slots not present in text_map: shape preserved (template text shows through).
    text_map keys not in slots: silently dropped (logs WARN).
    """
    slots = placeholder_map.get("slots") or []
    slot_by_id = {s["id"]: s for s in slots if isinstance(s, dict) and "id" in s}

    # Pyramid auto-dark for top-half tiers (only if slot doesn't already override)
    tier_keys = [k for k in text_map.keys() if k.startswith("tier_")]
    total_tiers = len(tier_keys)
    DARK_HEX = "#0B2A4A"

    # Two-pass: first replace texts, then remove shapes for slots whose new_text is empty.
    # (Removal must happen after replacement so tree_paths don't shift mid-loop.)
    shapes_to_blank_remove = []

    for slot_id, slot in slot_by_id.items():
        if slot_id not in text_map:
            continue
        new_text = text_map[slot_id]
        tree_path = slot.get("tree_path")
        if not tree_path:
            continue
        shape = _find_shape_by_tree_path(slide, tree_path)
        if shape is None or not shape.has_text_frame:
            print(f"  WARN tier1 mapping: slot {slot_id!r} tree_path={tree_path} 找不到 / 无 text_frame")
            continue
        # If new_text is empty AND slot allows removal (default True), queue for shape removal
        # (so leftover capsule/circle decoration vanishes instead of showing as ghost shape)
        if new_text == "" and slot.get("keep_when_empty") is not True:
            shapes_to_blank_remove.append(shape)
            continue
        # Color: slot override > auto-tier dark > none
        color = slot.get("text_color")
        if color is None and slot_id.startswith("tier_") and total_tiers > 0:
            try:
                tier_num = int(slot_id.split("_", 1)[1])
                if tier_num <= total_tiers // 2:
                    color = DARK_HEX
            except ValueError:
                pass
        font_size = slot.get("font_size_pt")
        _replace_shape_text(shape, new_text, text_color_hex=color,
                              font_size_pt=font_size)

    # Now remove the empty-slot shapes (top-level only via parent.remove)
    for shape in shapes_to_blank_remove:
        try:
            sp = shape._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)
        except Exception:
            # If shape is nested deeply inside groups, just blank the text (already done above for safety)
            try:
                _replace_shape_text(shape, "", text_color_hex=None, font_size_pt=None)
            except Exception:
                pass

    # Log any text_map keys that have no matching slot
    unmatched = set(text_map.keys()) - set(slot_by_id.keys())
    if unmatched:
        print(f"  WARN tier1 mapping: text_map keys {sorted(unmatched)} 无对应 slot, 已跳过")


def _apply_text_map_by_geometry(slide, text_map: dict[str, str]):
    """Legacy fallback: geometric (top, left) ordering of pattern-matched shapes."""
    placeholders = _collect_placeholder_shapes(slide)
    ordered_keys = sorted(text_map.keys(), key=lambda k: _key_sort_index(k))
    tier_keys = [k for k in ordered_keys if k.startswith("tier_")]
    total_tiers = len(tier_keys)
    DARK_HEX = "#0B2A4A"
    for idx, key in enumerate(ordered_keys):
        if idx >= len(placeholders):
            break
        color_override = None
        if key.startswith("tier_") and total_tiers > 0:
            try:
                tier_num = int(key.split("_", 1)[1])
                if tier_num <= total_tiers // 2:
                    color_override = DARK_HEX
            except ValueError:
                pass
        _replace_shape_text(placeholders[idx], text_map[key], text_color_hex=color_override)


def _key_sort_index(key: str) -> tuple:
    """Sort priority for text_map keys, aligned with author semantic intent.

    Order rules(matches placeholder geometric top-to-bottom in template):
      - Pure numeric → by value
      - With explicit index "_N" suffix(e.g. tier_3, card_2_title)→ by (prefix, N, sub-role)
      - sub-role priority: title < body < (others alphabetical) — title-first per author intent
      - Headings (number/title/header) precede content (body/desc/supporting)
    """
    import re
    # Pure numeric
    try:
        return (0, int(key), 0, "")
    except ValueError:
        pass
    # Structured: <prefix>_<num>[_<role>] like card_1_title / tier_2 / step_3_desc
    m = re.match(r"^([a-zA-Z]+)_(\d+)(?:_(.+))?$", key)
    if m:
        prefix, num, role = m.group(1), int(m.group(2)), m.group(3) or ""
        role_order = {"title": 0, "header": 0, "number": 0,
                      "body": 1, "desc": 1, "supporting": 1,
                      "subtitle": 2}.get(role, 3)
        return (1, prefix, num, role_order, role)
    # Sidebar / no-number keys (sidebar_left_title vs sidebar_left_body)
    m = re.match(r"^(.+?)_(title|header|body|desc|subtitle)$", key)
    if m:
        base, role = m.group(1), m.group(2)
        role_order = {"title": 0, "header": 0, "body": 1, "desc": 1, "subtitle": 2}.get(role, 3)
        return (2, base, 0, role_order, role)
    return (3, key, 0, 0, "")


def _load_template_prs_for_tier1(theme: ModuleType, plan_dir: str | None) -> Presentation | None:
    """Load template .pptx for tier1 reuse, if theme is a .pptx-extracted module."""
    # The template stem == theme name we passed (e.g. "template_golden")
    # Extract from theme.__name__ which is "extracted_<stem>"
    name = theme.__name__
    if not name.startswith("extracted_"):
        return None
    stem = name[len("extracted_"):]
    pptx_path = _find_template(stem, plan_dir)
    if pptx_path is None:
        return None
    return Presentation(str(pptx_path))


# ----- red_line_words 第 4 道防线 -----

def _parse_red_line_words(brief_path: str | Path) -> list[str]:
    """Parse brief.md frontmatter + yaml fences for constraints.red_line_words.

    Brief schema supports two yaml locations:
    - ``---\\n...yaml...\\n---`` frontmatter at top
    - any ````yaml`` fence in body (brainstorm writes the constraints block as a fence)

    Returns a list of word strings (empty if not found / brief unreadable).
    Non-fatal: callers decide whether missing words is an error.
    """
    import re
    try:
        import yaml as _yaml
    except ImportError:
        return []
    try:
        text = Path(brief_path).read_text(encoding="utf-8")
    except (FileNotFoundError, OSError):
        return []
    candidates: list[str] = []
    fm = re.match(r"^---\n(.*?)\n---", text, re.S)
    if fm:
        candidates.append(fm.group(1))
    candidates.extend(re.findall(r"```yaml\n(.*?)\n```", text, re.S))
    for block in candidates:
        try:
            data = _yaml.safe_load(block) or {}
        except Exception:
            continue
        if not isinstance(data, dict):
            continue
        constraints = data.get("constraints") or {}
        words = constraints.get("red_line_words") or []
        if isinstance(words, list) and words:
            return [str(w).strip() for w in words if str(w).strip()]
    return []


def _check_red_line_words(brief_path: str | Path | None,
                            content_md_path: str | Path | None,
                            deck_plan: dict[str, Any]) -> None:
    """第 4 道防线:author 自检 / critic Stage C/D 漏了也在 build 时拦。

    fail-loud raises ValueError. 跳过场景(全部静默):
    - brief_path 缺失 / 不可读
    - constraints.red_line_words 未配置或为空 list
    - content_md_path 缺失 / 不可读(只查 deck_plan)

    检查范围:
    - content.md 全文(若可读)
    - deck_plan.json 序列化后所有文本字段(含 builder Step 3 字数修复引入的新词)
    """
    if brief_path is None:
        return
    words = _parse_red_line_words(brief_path)
    if not words:
        return
    # content.md
    if content_md_path is not None:
        try:
            content_text = Path(content_md_path).read_text(encoding="utf-8")
        except (FileNotFoundError, OSError):
            content_text = ""
        for w in words:
            if w and w in content_text:
                raise ValueError(
                    f"红线词 {w!r} 在 content.md 残留 ({content_md_path}),"
                    f"critic 漏检 + author 自检也漏。fix: author rework 删该词。"
                )
    # deck_plan text fields(builder 自动修复 / 字数缩写时可能引入)
    plan_str = json.dumps(deck_plan, ensure_ascii=False)
    for w in words:
        if w and w in plan_str:
            raise ValueError(
                f"红线词 {w!r} 在 deck_plan.json 残留(可能 builder 自动修复时引入)。"
                f"fix: 检查最近一轮 builder 改写,改回不踩词的措辞。"
            )


# ----- build_deck -----

def build_deck(plan: dict[str, Any]) -> Path:
    """按 deck_plan 逐 slide 调 make_*,存 .pptx,返回输出路径。

    自动处理 3 个 cross-cutting 字段(build.py 集中负责,theme 不感知):
    - **footer**: 内容页(FOOTERED_LAYOUTS)统一加分隔线 + "N / TOTAL" + 可选元数据
    - **footer_meta**(plan 顶层): classification / project / version,显示在 footer 左侧
    - **source**(slide 级): 数据 slide 的引文,渲染在 footer 上方

    presentation_mode:plan 顶层可设 "speaker"(默认)或 "handout",
    会 set helpers.PRESENTATION_MODE,theme layout 据此切字号 / box 高度 / padding。
    """
    # set presentation mode(影响 layout 字数 / 字号)
    mode = plan.get("presentation_mode", "speaker")
    if mode not in ("speaker", "handout"):
        raise ValueError(f"presentation_mode 必须是 speaker / handout,得到 {mode!r}")
    H.PRESENTATION_MODE = mode

    # 第 4 道防线:lock content 前 grep 红线词(brief / content / deck_plan)
    # brief_path + content_md_path 由主线程派 builder 时 inject(deck_plan.json 顶层)。
    # 都缺则跳过(本地手写 plan 跑 build.py 不强制有 brief)。
    _check_red_line_words(
        plan.get("brief_path"),
        plan.get("content_md_path"),
        plan,
    )

    theme = load_theme(plan["theme"], plan.get("_plan_dir"))

    # Tier1 模板复用:若任一 slide 有 `tier1_template_page` 字段,以模板 prs 作起点
    # (保留模板 master + layouts + theme,让复制过去的 shapes 视觉对齐),
    # 先删除模板自带的 38 张 slide,再按 plan 顺序加 slide。
    use_tier1 = any("tier1_template_page" in s for s in plan["slides"])
    source_prs: Presentation | None = None
    if use_tier1:
        source_prs = _load_template_prs_for_tier1(theme, plan.get("_plan_dir"))
        if source_prs is None:
            raise ValueError(
                f"tier1 path 启用,但找不到 theme {plan['theme']!r} 对应的 .pptx 模板。"
                "tier1 复用要求 theme 是 ingested 模板(library/pptx-templates/_source/<name>.pptx)。"
            )
        # 以模板为起点(保留 theme schemes + masters)
        template_path = _find_template(plan["theme"], plan.get("_plan_dir"))
        prs = Presentation(str(template_path))
        prs.slide_width = H.SLIDE_W
        prs.slide_height = H.SLIDE_H
        # 删除模板原所有 slide(必须 drop_rel + 删 sldIdLst entry,否则 part 残留
        # 导致 LibreOffice 看到 [Content_Types].xml 声明的 slide 文件但无引用,加载失败)
        xml_slides = prs.slides._sldIdLst
        slide_id_lst = list(xml_slides)
        for sld_id in slide_id_lst:
            rId = sld_id.rId
            prs.part.drop_rel(rId)
            xml_slides.remove(sld_id)
    else:
        prs = Presentation()
        prs.slide_width = H.SLIDE_W
        prs.slide_height = H.SLIDE_H

    footer_meta = plan.get("footer_meta", {}) or {}

    # 预扫:算 footer 页总数(用于 "N / TOTAL" 的 TOTAL)
    total_footered = sum(
        1 for s in plan["slides"] if s["layout"] in FOOTERED_LAYOUTS
    )
    footer_idx = 0

    for i, slide in enumerate(plan["slides"], 1):
        layout = slide["layout"]
        # tier1 path:从 source_prs 复制指定 page 的 slide
        if "tier1_template_page" in slide:
            if source_prs is None:
                raise ValueError(f"第 {i} 页用 tier1 但 source_prs 未加载")
            page_idx = slide["tier1_template_page"]
            if not (0 <= page_idx < len(source_prs.slides)):
                raise ValueError(
                    f"第 {i} 页 tier1_template_page={page_idx} 越界 "
                    f"(模板有 {len(source_prs.slides)} 页)"
                )
            src_slide = source_prs.slides[page_idx]
            new_slide = _copy_slide_from_source(src_slide, prs)
            # 加载 placeholder_map(若 library/pptx-templates/items/.../<page>/placeholder_map.yaml 存在)
            # 优先级:slide 显式 placeholder_map > 自动按 tier1_template_page 查找
            pmap = slide.get("placeholder_map")
            if pmap is None:
                pmap = _load_placeholder_map(theme, page_idx, plan.get("_plan_dir"))
            # 替换 title + placeholder text + 移除 iSlide footer
            _apply_tier1_text_map(
                new_slide,
                title=slide.get("title"),
                text_map=slide.get("text_map", {}),
                placeholder_map=pmap,
            )
            # source 引文 + footer 用 deck 自己的系统
            source = slide.get("source")
            if source:
                H.source_citation(new_slide, source)
            if layout in FOOTERED_LAYOUTS and total_footered > 0:
                footer_idx += 1
                H.footer(
                    new_slide, footer_idx, total_footered,
                    classification=footer_meta.get("classification"),
                    project=footer_meta.get("project"),
                    version=footer_meta.get("version"),
                )
            continue
        # tier2 path:Python theme make_*
        fn = getattr(theme, f"make_{layout}", None)
        if fn is None:
            # Fail loud(不 silent remap)— 让 author / critic / 用户立刻看到错配,
            # 不让 builder 默默 fallback 到 bullet_list 等掩盖,等 audience 才暴露
            available = sorted(name[len("make_"):] for name in dir(theme)
                                if name.startswith("make_") and callable(getattr(theme, name)))
            theme_stem = theme.__name__.replace("extracted_", "")
            raise ValueError(
                f"第 {i} 页 layout={layout!r}: theme {theme.__name__!r} 无 make_{layout}。\n"
                f"可能原因:选了 tier1-only 模板(meta.yaml.tier2_python_theme=null),"
                f"但 deck_plan 这页没用 tier1 路径。\n"
                f"3 个 fix 选项:\n"
                f"  ① 让 author 改 layout 到 theme 支持清单:{available}\n"
                f"  ② 启用 tier1 路径:deck_plan slide 加 tier1_template_page + text_map "
                f"(参考 library/pptx-templates/items/{theme_stem}/pages/*/placeholder_map.yaml)\n"
                f"  ③ 主线程实现 themes/{theme_stem}.py 的 make_{layout} 函数"
            )
        # 弹出 cross-cutting 字段,不传给 make_* fn(避免 TypeError)
        fields = {k: v for k, v in slide.items() if k != "layout"}
        source = fields.pop("source", None)
        try:
            fn(prs, **fields)
        except TypeError as e:
            raise ValueError(f"第 {i} 页 layout={layout}: {e}") from e

        # 数据引文(source)→ footer 上方
        if source:
            H.source_citation(prs.slides[-1], source)

        # 页脚 + 页码 + footer_meta
        if layout in FOOTERED_LAYOUTS and total_footered > 0:
            footer_idx += 1
            H.footer(
                prs.slides[-1], footer_idx, total_footered,
                classification=footer_meta.get("classification"),
                project=footer_meta.get("project"),
                version=footer_meta.get("version"),
            )

    out = Path(plan["output"]).expanduser()
    if not out.is_absolute():
        out = (Path(plan["_plan_dir"]) / out).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


# ----- render -----

def render(pptx_path: str | Path, out_dir: str | Path) -> list[Path]:
    """soffice → PDF → pdftoppm → 逐页 PNG。返回 PNG 路径列表。"""
    if shutil.which("soffice") is None:
        raise RuntimeError("soffice 未安装。请: brew install --cask libreoffice")
    if shutil.which("pdftoppm") is None:
        raise RuntimeError("pdftoppm 未安装。请: brew install poppler")
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             str(pptx_path), "--outdir", str(out_dir)],
            check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"soffice 转 PDF 失败: {e.stderr}") from e
    pdf = out_dir / (pptx_path.stem + ".pdf")
    if not pdf.exists():
        raise RuntimeError(f"soffice 跑了但未产 PDF: {pdf}")
    try:
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "120", str(pdf), str(out_dir / "page")],
            check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"pdftoppm 转 PNG 失败: {e.stderr}") from e
    return sorted(out_dir.glob("page-*.jpg"))


# ----- CLI -----

def main(argv: list[str]) -> None:
    if not argv:
        sys.exit("用法: python3 build.py deck_plan.json [--no-render]")
    plan_path = argv[0]
    do_render = "--no-render" not in argv
    plan = load_plan(plan_path)
    out = build_deck(plan)
    print(f"已生成 {out}")
    if do_render:
        render_dir = out.parent / (out.stem + "_render")
        pngs = render(out, render_dir)
        print(f"已渲染 {len(pngs)} 页 → {render_dir}")


if __name__ == "__main__":
    main(sys.argv[1:])
