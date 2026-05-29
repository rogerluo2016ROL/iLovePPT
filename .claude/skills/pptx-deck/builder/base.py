"""builder/base.py —— 共享层:plan/theme load + build_deck orchestrator + render。

负责:
- load_plan / load_theme(单 theme · 旧 API,向后兼容)
- ThemeSpec + parse_theme(P3-9 多模板组合 deck 支持 str | list | dict)
- _extract_design_tokens / _extract_theme_from_pptx(.pptx → theme module)
- _find_template / _list_available_templates(模板查找)
- _check_red_line_words / _parse_red_line_words(第 4 道防线)
- build_deck(主 orchestrator,逐 slide 派 tier1 / tier2 / tier3)
- render(soffice → PDF → pdftoppm → JPG)

注意:tier1 / tier2 / tier3 派发只在 build_deck 内,不在本模块定义 tier 实现。
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path
from types import ModuleType
from typing import Any

from pptx import Presentation

# === P0-3 · 静默吞错可见性(组件 A) ===
# build 过程中被"吞掉但回落"的错误统一走这里:既保留回落鲁棒性,又不再静默。
# 模块级列表方便测试断言;build_deck 入口清空,结束打印汇总。
BUILD_WARNINGS: list[str] = []


def _warn(stage: str, msg: str) -> None:
    """记一条 build warning:append 到 BUILD_WARNINGS + 打印到 stderr。

    stage 约定前缀:`builder.token-extract` / `builder.red-line` /
    `tier1.slot-map` / `tier1.shape-removal`。
    """
    line = f"[{stage}] WARN {msg}"
    BUILD_WARNINGS.append(line)
    print(line, file=sys.stderr)


HERE = Path(__file__).resolve().parent  # <repo>/.claude/skills/pptx-deck/builder
PPTX_DECK_DIR = HERE.parent              # <repo>/.claude/skills/pptx-deck
PPTX_SKILL_DIR = PPTX_DECK_DIR.parent / "pptx"  # <repo>/.claude/skills/pptx

# 把 pptx-deck + pptx 加入 sys.path(让 `import helpers` / `from themes import ...` 可用)。
# 幂等,重复加无害。
for _p in [str(PPTX_SKILL_DIR), str(PPTX_DECK_DIR)]:
    if _p not in sys.path:
        sys.path.insert(0, _p)

import helpers as H
from themes import tech_blue as _tech_blue
from themes import template_golden as _template_golden
from themes import template_training as _template_training

THEMES: dict[str, ModuleType] = {
    "tech_blue": _tech_blue,
    "template_training": _template_training,
}

# Themes with Python layout implementations but **not** registered in THEMES.
# These are ingested .pptx templates (library/pptx-templates/items/<name>/) whose
# theme module supplies layout shapes; color/font tokens come from the .pptx at
# runtime via _extract_theme_from_pptx. Keying off the .pptx stem lets
# `_extract_theme_from_pptx` pick the right base module instead of always
# hardcoding tech_blue.
PPTX_BASE_THEMES: dict[str, ModuleType] = {
    "template_golden": _template_golden,
}

# 需要页脚 + 页码的 layout(规范:visual-qa.md §页脚 / 页码完整性)。
# cover / section_divider / closing 不计入页码。
FOOTERED_LAYOUTS: frozenset[str] = frozenset({
    "toc", "single_focus", "compare", "compare_pk", "matrix_2x2", "cards",
    "bullet_list", "table", "pic_text", "summary",
    "timeline_band_3", "tri_pyramid_4sub_3", "cards_flag_3",
})


# ===========================================================================
# ThemeSpec + parse_theme(P3-9 多模板组合 deck)
# ===========================================================================

@dataclass
class ThemeSpec:
    """多模板组合 deck 的 theme 配置,解析自 brief.theme(str | list | dict)。

    向后兼容:
    - 单 str(legacy) → ThemeSpec(default=<str>)
    - 单 list[1] → ThemeSpec(default=<list[0]>)
    - list[2+] → ThemeSpec(default=<list[0]>, overrides 留空,需 deck_plan 显式或后续策略)
    - dict {default, overrides} → ThemeSpec(default=..., overrides=...)

    overrides schema(P3-9):
    - key 是页号(1-indexed,字符串或 int)或范围 "N-M"
    - value 是 theme id(str)
    - resolve_for_page(idx, total) 把范围展开,返回该页该用的 theme id

    跨模板视觉协调(字体/色板)留 TODO,本任务保模板 1 不混搭。
    """
    default: str
    overrides: dict[str, str] = field(default_factory=dict)
    # extras: list[str] 字段保留多 list 模式下声明的"备选模板",当前 build 路径未消费,
    # 留给 deck_plan 显式 slide.theme_override 或后续 chapter override 策略用。
    extras: list[str] = field(default_factory=list)

    def all_themes(self) -> list[str]:
        """返回去重后的所有 theme id(default + overrides 各值 + extras)。"""
        seen = {self.default}
        out = [self.default]
        for v in list(self.overrides.values()) + self.extras:
            if v and v not in seen:
                seen.add(v)
                out.append(v)
        return out

    def resolve_for_page(self, page_idx_1based: int) -> str:
        """返回第 N 页(1-indexed)该用的 theme id。

        匹配优先级:
        1. overrides 中精确 page number 命中(key="5" 或 5)
        2. overrides 中范围命中(key="5-8" 表示第 5-8 页)
        3. fallback 到 default

        范围解析:"N-M" → [N, M] 闭区间;非法 key 静默忽略(返回 default)。
        """
        # 1. 精确匹配
        for key in (str(page_idx_1based), page_idx_1based):
            if key in self.overrides:
                return self.overrides[key]
        # 2. 范围匹配
        for key, value in self.overrides.items():
            key_s = str(key)
            if "-" in key_s:
                parts = key_s.split("-", 1)
                try:
                    lo, hi = int(parts[0].strip()), int(parts[1].strip())
                except ValueError:
                    continue
                if lo <= page_idx_1based <= hi:
                    return value
        return self.default


def parse_theme(brief_theme: Any) -> ThemeSpec:
    """把 brief.theme(str | list | dict)规范化成 ThemeSpec。

    支持的 schema:
    - **str**(legacy):`"enterprise_skyline"` → ThemeSpec(default="enterprise_skyline")
    - **list[str]**:`["enterprise_skyline", "finance_arrow"]` → ThemeSpec(default="enterprise_skyline", extras=["finance_arrow"])
    - **dict**:
        ```
        {
            "default": "enterprise_skyline",
            "overrides": {"1": "enterprise_skyline", "5-8": "finance_arrow"}
        }
        ```
        → ThemeSpec(default="enterprise_skyline", overrides={"1": ..., "5-8": ...})

    错误:
    - None / 空 list / 空 dict → ValueError
    - 非 str/list/dict → ValueError
    - list[0] 不是 str → ValueError
    - dict 缺 default 字段 → ValueError
    """
    if brief_theme is None or brief_theme == "" or brief_theme == [] or brief_theme == {}:
        raise ValueError(f"brief.theme 不能为空: {brief_theme!r}")

    # 单 str(legacy · 最常见路径)
    if isinstance(brief_theme, str):
        return ThemeSpec(default=brief_theme)

    # list 形式
    if isinstance(brief_theme, list):
        if not all(isinstance(x, str) for x in brief_theme):
            raise ValueError(f"brief.theme list 所有元素必须是 str: {brief_theme!r}")
        if len(brief_theme) == 0:
            raise ValueError("brief.theme list 为空")
        default = brief_theme[0]
        extras = list(brief_theme[1:])
        return ThemeSpec(default=default, extras=extras)

    # dict 形式
    if isinstance(brief_theme, dict):
        default = brief_theme.get("default")
        if not default or not isinstance(default, str):
            raise ValueError(f"brief.theme dict 必须含 default(str)字段: {brief_theme!r}")
        overrides_raw = brief_theme.get("overrides") or {}
        if not isinstance(overrides_raw, dict):
            raise ValueError(f"brief.theme.overrides 必须是 dict: {overrides_raw!r}")
        # 全部 key 转 str(允许 yaml int key 例如 1)
        overrides = {str(k): str(v) for k, v in overrides_raw.items() if v}
        extras_raw = brief_theme.get("extras") or []
        extras = [str(x) for x in extras_raw if x]
        return ThemeSpec(default=default, overrides=overrides, extras=extras)

    raise ValueError(
        f"brief.theme 必须是 str / list[str] / dict{{default,overrides}},"
        f"得到 {type(brief_theme).__name__}: {brief_theme!r}"
    )


# ===========================================================================
# load_plan
# ===========================================================================

def load_plan(path: str | Path) -> dict[str, Any]:
    """读 + 校验 deck_plan.json。记录 _plan_dir 供相对 output 解析。"""
    p = Path(path)
    data = json.loads(p.read_text(encoding="utf-8"))
    for field_name in ("theme", "output", "slides"):
        if field_name not in data:
            raise ValueError(f"deck_plan 缺字段: {field_name}")
    if not isinstance(data["slides"], list) or not data["slides"]:
        raise ValueError("deck_plan.slides 必须是非空 list")
    for i, slide in enumerate(data["slides"], 1):
        if "layout" not in slide:
            raise ValueError(f"deck_plan 第 {i} 页缺 layout 字段")
    data["_plan_dir"] = str(p.resolve().parent)
    return data


# ===========================================================================
# theme load + .pptx token 提取
# ===========================================================================

def _extract_design_tokens(pptx_path: str) -> dict[str, Any]:
    """从 .pptx 提取扩展 design token(accent1-6 + 字号阶梯).

    返回字段(全 best-effort,缺则不在 dict 里):
    - font_header / font_body: master ea typeface
    - primary: RGBColor(accent1)
    - accent_2/3/4/5/6: RGBColor(accent2-6)
    - dk1 / lt1: RGBColor(主文本色 / 主背景色)
    - title_size_pt / body_size_pt: master 默认字号(int pt)
    """
    from lxml import etree
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    tokens: dict[str, Any] = {}

    def _hex2rgb(hx: str) -> RGBColor | None:
        if len(hx) == 6:
            return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        return None

    try:
        prs = Presentation(pptx_path)
    except Exception:
        return tokens

    # === master ea typeface + 字号 ===
    try:
        if prs.slide_masters:
            done = False
            for ph in prs.slide_masters[0].placeholders:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(qn("a:rPr"))
                        if rPr is not None:
                            ea = rPr.find(qn("a:ea"))
                            if ea is not None and ea.get("typeface"):
                                tokens["font_header"] = ea.get("typeface")
                                tokens["font_body"] = ea.get("typeface")
                                done = True
                                break
                    if done:
                        break
                if done:
                    break
    except Exception:
        pass

    # === 从 master XML 抽字号(直接读 slideMaster1.xml) ===
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
              "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        for part in prs.part.package.iter_parts():
            pn = part.partname
            if "slideMaster" in pn and pn.endswith(".xml"):
                root = etree.fromstring(part.blob)
                title_def = root.find(".//p:titleStyle//a:lvl1pPr/a:defRPr", ns)
                if title_def is not None:
                    sz = title_def.get("sz")
                    if sz and sz.isdigit():
                        tokens["title_size_pt"] = int(sz) // 100
                body_def = root.find(".//p:bodyStyle//a:lvl1pPr/a:defRPr", ns)
                if body_def is not None:
                    sz = body_def.get("sz")
                    if sz and sz.isdigit():
                        tokens["body_size_pt"] = int(sz) // 100
                break
    except Exception:
        pass

    # === theme1.xml: accent1-6 + dk1/lt1 ===
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for part in prs.part.package.iter_parts():
            pn = part.partname
            if "theme" in pn and pn.endswith(".xml"):
                root = etree.fromstring(part.blob)
                scheme = root.find(".//a:clrScheme", ns)
                if scheme is None:
                    continue
                for tag, key in [("accent1", "primary"),
                                  ("accent2", "accent_2"),
                                  ("accent3", "accent_3"),
                                  ("accent4", "accent_4"),
                                  ("accent5", "accent_5"),
                                  ("accent6", "accent_6"),
                                  ("dk1", "dk1"), ("lt1", "lt1")]:
                    node = scheme.find(f"a:{tag}", ns)
                    if node is None:
                        continue
                    srgb = node.find(".//a:srgbClr", ns)
                    if srgb is not None:
                        rgb = _hex2rgb(srgb.get("val", ""))
                        if rgb:
                            tokens[key] = rgb
                            continue
                    # dk1/lt1 可能是 sysClr,取 lastClr
                    if tag in ("dk1", "lt1"):
                        sys_clr = node.find("a:sysClr", ns)
                        if sys_clr is not None:
                            rgb = _hex2rgb(sys_clr.get("lastClr", ""))
                            if rgb:
                                tokens[key] = rgb
                break
    except Exception:
        pass

    return tokens


def _extract_theme_from_pptx(pptx_path: str) -> ModuleType:
    """从用户 .pptx 提取主色与字体,派生临时主题模块。

    Base 选择:用 .pptx 文件 stem 查 PPTX_BASE_THEMES(如 template_golden 用其
    专属 module 提供独有 layout),未注册则回落到 tech_blue。再从 base 源码加载
    全新 module 实例,用提取的 token 覆盖 FONT_* / PRIMARY 等属性。
    token 提取 best-effort,未提取到的保留 base 默认。
    """
    import importlib.util
    tokens = _extract_design_tokens(pptx_path)
    stem = Path(pptx_path).stem
    base_module = PPTX_BASE_THEMES.get(stem, _tech_blue)
    base_path = Path(base_module.__file__)
    out_name = f"extracted_{stem}"
    spec = importlib.util.spec_from_file_location(out_name, base_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法从 {base_path} 加载主题")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    # 字体
    if "font_header" in tokens:
        mod.FONT_HEADER = tokens["font_header"]
        mod.FONT_BODY = tokens.get("font_body", tokens["font_header"])
    # 主色 + 次级色(accent2-6 + dk1/lt1)
    if "primary" in tokens:
        mod.PRIMARY = tokens["primary"]
    for key in ("accent_2", "accent_3", "accent_4", "accent_5", "accent_6",
                "dk1", "lt1"):
        if key in tokens:
            setattr(mod, key.upper(), tokens[key])
    # 字号阶梯(若 master 显式定义)
    if "title_size_pt" in tokens:
        mod.TITLE_SIZE_PT = tokens["title_size_pt"]
    if "body_size_pt" in tokens:
        mod.BODY_SIZE_PT = tokens["body_size_pt"]

    font_status = tokens.get("font_header", "默认 Microsoft YaHei")
    color_status = tokens.get("primary", "默认 tech_blue 主色")
    print(f"  从模板提取主题: {out_name}")
    print(f"     字体: {font_status}")
    print(f"     主色: {color_status}")
    extra_accents = [k for k in ("accent_2", "accent_3", "accent_4", "accent_5",
                                  "accent_6") if k in tokens]
    if extra_accents:
        print(f"     次级色: {', '.join(extra_accents)}")
    if "title_size_pt" in tokens or "body_size_pt" in tokens:
        print(f"     字号阶梯: title={tokens.get('title_size_pt', '-')}pt /"
              f" body={tokens.get('body_size_pt', '-')}pt")
    return mod


def _repo_templates_dir() -> Path:
    """仓库根的 library/pptx-templates/_source/ 目录(模板 .pptx 全局共享)。
    base.py 位于 <repo>/.claude/skills/pptx-deck/builder/base.py。"""
    repo_root = Path(__file__).resolve().parents[4]  # base.py → builder → pptx-deck → skills → .claude → repo
    return repo_root / "library" / "pptx-templates" / "_source"


def _find_template(name: str, plan_dir: str | None = None) -> Path | None:
    """按短名查找 .pptx 模板。

    优先级:
      1. <plan_dir>/templates/<name>.pptx                       (deck 项目专属, 向后兼容)
      2. <repo>/library/pptx-templates/_source/<name>.pptx      (全局共享, 新位置)

    找到返回 Path,找不到返回 None。
    """
    candidates: list[Path] = []
    if plan_dir:
        candidates.append(Path(plan_dir) / "templates" / f"{name}.pptx")
    candidates.append(_repo_templates_dir() / f"{name}.pptx")
    for p in candidates:
        if p.exists():
            return p
    return None


def _list_available_templates() -> list[str]:
    """返回 <repo>/library/pptx-templates/_source/ 下所有 .pptx 短名(不含扩展名)。"""
    tdir = _repo_templates_dir()
    if not tdir.exists():
        return []
    return sorted(p.stem for p in tdir.glob("*.pptx"))


def load_theme(theme_id: str, plan_dir: str | None = None) -> ModuleType:
    """解析 theme_id 到 theme 模块。

    Args:
        theme_id: 三种形式之一
            - 内置 theme 名(如 "tech_blue")
            - 短名(如 "company_a")—— 查找 library/pptx-templates/_source/<name>.pptx(全局)
              或 <plan_dir>/templates/<name>.pptx(deck 项目本地)
            - 路径(含 "/" 或以 ".pptx" 结尾)—— 直接当 .pptx 路径
        plan_dir: deck plan 所在目录,影响相对路径解析 + 短名查找优先级
    """
    if theme_id in THEMES:
        return THEMES[theme_id]
    # 含 / 或以 .pptx 结尾 → 当路径处理
    if str(theme_id).endswith(".pptx") or "/" in str(theme_id):
        path = Path(theme_id).expanduser()
        if not path.is_absolute() and plan_dir:
            path = (Path(plan_dir) / path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"theme .pptx 不存在: {path}")
        return _extract_theme_from_pptx(str(path))
    # 短名 → 查 library/pptx-templates/_source/ + <plan_dir>/templates/
    found = _find_template(theme_id, plan_dir)
    if found is not None:
        return _extract_theme_from_pptx(str(found))
    # 未找到 → 列可用的帮用户排错
    available = _list_available_templates()
    available_str = ", ".join(available) if available else "(空,把 .pptx 放进 library/pptx-templates/_source/)"
    raise ValueError(
        f"未知 theme: {theme_id!r}. "
        f"内置: tech_blue. "
        f"library/pptx-templates/_source/ 可用: {available_str}. "
        f"或直接给 .pptx 绝对/相对路径。"
    )


# ===========================================================================
# red_line_words 第 4 道防线
# ===========================================================================

def _parse_red_line_words(brief_path: str | Path) -> list[str]:
    """Parse brief.md frontmatter + yaml fences for constraints.red_line_words.

    Brief schema supports two yaml locations:
    - ``---\\n...yaml...\\n---`` frontmatter at top
    - any ````yaml`` fence in body (brainstorm writes the constraints block as a fence)

    Returns a list of word strings (empty if not found / brief unreadable).
    Non-fatal: callers decide whether missing words is an error.
    """
    import re
    try:
        import yaml as _yaml
    except ImportError:
        return []
    try:
        text = Path(brief_path).read_text(encoding="utf-8")
    except (FileNotFoundError, OSError):
        return []
    candidates: list[str] = []
    fm = re.match(r"^---\n(.*?)\n---", text, re.S)
    if fm:
        candidates.append(fm.group(1))
    candidates.extend(re.findall(r"```yaml\n(.*?)\n```", text, re.S))
    for block in candidates:
        try:
            data = _yaml.safe_load(block) or {}
        except Exception:
            continue
        if not isinstance(data, dict):
            continue
        constraints = data.get("constraints") or {}
        words = constraints.get("red_line_words") or []
        if isinstance(words, list) and words:
            return [str(w).strip() for w in words if str(w).strip()]
    return []


def _check_red_line_words(brief_path: str | Path | None,
                            content_md_path: str | Path | None,
                            deck_plan: dict[str, Any]) -> None:
    """第 4 道防线:author 自检 / critic Stage C/D 漏了也在 build 时拦。

    fail-loud raises ValueError. 跳过场景(全部静默):
    - brief_path 缺失 / 不可读
    - constraints.red_line_words 未配置或为空 list
    - content_md_path 缺失 / 不可读(只查 deck_plan)

    检查范围:
    - content.md 全文(若可读)
    - deck_plan.json 序列化后所有文本字段(含 builder Step 3 字数修复引入的新词)
    """
    if brief_path is None:
        return
    words = _parse_red_line_words(brief_path)
    if not words:
        return
    # content.md
    if content_md_path is not None:
        try:
            content_text = Path(content_md_path).read_text(encoding="utf-8")
        except (FileNotFoundError, OSError):
            content_text = ""
        for w in words:
            if w and w in content_text:
                raise ValueError(
                    f"红线词 {w!r} 在 content.md 残留 ({content_md_path}),"
                    f"critic 漏检 + author 自检也漏。fix: author rework 删该词。"
                )
    # deck_plan text fields(builder 自动修复 / 字数缩写时可能引入)
    plan_str = json.dumps(deck_plan, ensure_ascii=False)
    for w in words:
        if w and w in plan_str:
            raise ValueError(
                f"红线词 {w!r} 在 deck_plan.json 残留(可能 builder 自动修复时引入)。"
                f"fix: 检查最近一轮 builder 改写,改回不踩词的措辞。"
            )


# ===========================================================================
# build_deck —— orchestrator(逐 slide 派 tier1 / tier2 / tier3)
# ===========================================================================

def build_deck(plan: dict[str, Any]) -> Path:
    """按 deck_plan 逐 slide 调 make_*,存 .pptx,返回输出路径。

    自动处理 3 个 cross-cutting 字段(build.py 集中负责,theme 不感知):
    - **footer**: 内容页(FOOTERED_LAYOUTS)统一加分隔线 + "N / TOTAL" + 可选元数据
    - **footer_meta**(plan 顶层): classification / project / version,显示在 footer 左侧
    - **source**(slide 级): 数据 slide 的引文,渲染在 footer 上方

    presentation_mode:plan 顶层可设 "speaker"(默认)或 "handout",
    会 set helpers.PRESENTATION_MODE,theme layout 据此切字号 / box 高度 / padding。

    Theme 解析(P3-9 多模板组合):plan["theme"] 可以是 str / list / dict;
    parse_theme() 规范化成 ThemeSpec。当前 tier1 / tier2 path 走 ThemeSpec.default 模板
    的色板 / 字体(避免跨模板字体色板混搭),后续可扩展按 slide 选 theme override。
    """
    # 延迟 import 避免循环依赖(tier1/tier2/tier3 都依赖 base.py 的常量)
    from . import tier1, tier2, tier3
    BUILD_WARNINGS.clear()

    # set presentation mode(影响 layout 字数 / 字号)
    mode = plan.get("presentation_mode", "speaker")
    if mode not in ("speaker", "handout"):
        raise ValueError(f"presentation_mode 必须是 speaker / handout,得到 {mode!r}")
    H.PRESENTATION_MODE = mode

    # 第 4 道防线:lock content 前 grep 红线词(brief / content / deck_plan)
    _check_red_line_words(
        plan.get("brief_path"),
        plan.get("content_md_path"),
        plan,
    )

    # P3-9:解析 theme(str | list | dict → ThemeSpec)。
    # 当前 build path 只用 ThemeSpec.default 的模板做主 prs init(保模板 1 不混搭色板字体)。
    # ThemeSpec.overrides / extras 后续可扩展用作 chapter override;本任务只做"机械可工作"路径。
    theme_spec = parse_theme(plan["theme"])
    theme = load_theme(theme_spec.default, plan.get("_plan_dir"))

    # 把 theme_spec 写回 plan 便于 tier1/tier2 内部 inspect(诊断用,不改 plan 语义)
    plan["_theme_spec"] = theme_spec

    # Tier1 模板复用:若任一 slide 有 `tier1_template_page` 字段,以模板 prs 作起点
    use_tier1 = any("tier1_template_page" in s for s in plan["slides"])
    source_prs: Presentation | None = None
    if use_tier1:
        source_prs = tier1.load_template_prs(theme, plan.get("_plan_dir"))
        if source_prs is None:
            raise ValueError(
                f"tier1 path 启用,但找不到 theme {theme_spec.default!r} 对应的 .pptx 模板。"
                "tier1 复用要求 theme 是 ingested 模板(library/pptx-templates/_source/<name>.pptx)。"
            )
        prs = tier1.init_prs_from_template(theme_spec.default, plan.get("_plan_dir"))
    else:
        prs = Presentation()
        prs.slide_width = H.SLIDE_W
        prs.slide_height = H.SLIDE_H

    footer_meta = plan.get("footer_meta", {}) or {}

    # 预扫:算 footer 页总数(用于 "N / TOTAL" 的 TOTAL)
    total_footered = sum(
        1 for s in plan["slides"] if s["layout"] in FOOTERED_LAYOUTS
    )
    footer_idx = 0

    for i, slide in enumerate(plan["slides"], 1):
        layout = slide["layout"]
        # tier1 path:从 source_prs 复制指定 page 的 slide
        if "tier1_template_page" in slide:
            if source_prs is None:
                raise ValueError(f"第 {i} 页用 tier1 但 source_prs 未加载")
            tier1.render_tier1_slide(
                prs=prs,
                source_prs=source_prs,
                slide_def=slide,
                page_no=i,
                theme=theme,
                plan_dir=plan.get("_plan_dir"),
            )
            # source 引文 + footer 用 deck 自己的系统
            source = slide.get("source")
            if source:
                H.source_citation(prs.slides[-1], source)
            if layout in FOOTERED_LAYOUTS and total_footered > 0:
                footer_idx += 1
                H.footer(
                    prs.slides[-1], footer_idx, total_footered,
                    classification=footer_meta.get("classification"),
                    project=footer_meta.get("project"),
                    version=footer_meta.get("version"),
                )
            continue
        # tier2 path:Python theme make_*
        try:
            tier2.render_tier2_slide(prs=prs, theme=theme, slide_def=slide, page_no=i)
        except tier2.LayoutNotFoundError as e:
            # tier3 fallback / 错误兜底:目前是 fail-loud(不 silent remap)
            tier3.handle_missing_layout(theme=theme, slide_def=slide, page_no=i, original_err=e)

        # 数据引文(source)→ footer 上方
        source = slide.get("source")
        if source:
            H.source_citation(prs.slides[-1], source)

        # 页脚 + 页码 + footer_meta
        if layout in FOOTERED_LAYOUTS and total_footered > 0:
            footer_idx += 1
            H.footer(
                prs.slides[-1], footer_idx, total_footered,
                classification=footer_meta.get("classification"),
                project=footer_meta.get("project"),
                version=footer_meta.get("version"),
            )

    out = Path(plan["output"]).expanduser()
    if not out.is_absolute():
        out = (Path(plan["_plan_dir"]) / out).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    if BUILD_WARNINGS:
        print(f"[build] {len(BUILD_WARNINGS)} warnings — {out}", file=sys.stderr)
    return out


# ===========================================================================
# render
# ===========================================================================

def render(pptx_path: str | Path, out_dir: str | Path) -> list[Path]:
    """soffice → PDF → pdftoppm → 逐页 PNG。返回 PNG 路径列表。"""
    if shutil.which("soffice") is None:
        raise RuntimeError("soffice 未安装。请: brew install --cask libreoffice")
    if shutil.which("pdftoppm") is None:
        raise RuntimeError("pdftoppm 未安装。请: brew install poppler")
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             str(pptx_path), "--outdir", str(out_dir)],
            check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"soffice 转 PDF 失败: {e.stderr}") from e
    pdf = out_dir / (pptx_path.stem + ".pdf")
    if not pdf.exists():
        raise RuntimeError(f"soffice 跑了但未产 PDF: {pdf}")
    try:
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "120", str(pdf), str(out_dir / "page")],
            check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"pdftoppm 转 PNG 失败: {e.stderr}") from e
    return sorted(out_dir.glob("page-*.jpg"))
