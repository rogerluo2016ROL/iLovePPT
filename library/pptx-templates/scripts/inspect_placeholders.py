"""Generate placeholder_map.yaml.draft skeleton · script computes tree_path, LLM fills semantics.

Usage:
    inspect_placeholders.py <pptx_path> <page_idx>      # page_idx 0-indexed

Output (stdout, YAML):
    template_page_index: 0
    layout_class: "?"          # extractor sets to match meta.yaml.draft.layout_type
    slots:
      - id: "?"                # extractor fills with "title" / "tier_1" / "card_1_body"
        tree_path: '3'
        raw_text_sample: "深蓝主标题"
        bbox: { left: 1.0, top: 2.0, width: 8.0, height: 1.0 }   # inches
        font_size_pt: 44
        capacity_chars: "?"    # extractor estimates
        text_color_override: null

Exit codes:
  0 = OK
  1 = pptx not found / page out of range
  2 = pptx parse error
"""
from __future__ import annotations

import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu


def walk_shapes(shapes, prefix: str = "") -> list[dict]:
    """Recursively walk shapes, return all leaf shapes containing text."""
    out = []
    for i, shape in enumerate(shapes):
        path = f"{prefix}{i}" if not prefix else f"{prefix}.{i}"
        if shape.shape_type == 6:  # group
            try:
                out.extend(walk_shapes(shape.shapes, prefix=path))
            except (AttributeError, ValueError):
                pass
        else:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = tf.text.strip() if tf.text else ""
            if not text:
                continue
            font_size = None
            try:
                run = tf.paragraphs[0].runs[0]
                if run.font.size:
                    font_size = run.font.size.pt
            except (IndexError, AttributeError):
                pass
            out.append({
                "tree_path": path,
                "raw_text_sample": text[:60],
                "bbox": {
                    "left": round(Emu(shape.left).inches, 2) if shape.left is not None else 0.0,
                    "top": round(Emu(shape.top).inches, 2) if shape.top is not None else 0.0,
                    "width": round(Emu(shape.width).inches, 2) if shape.width is not None else 0.0,
                    "height": round(Emu(shape.height).inches, 2) if shape.height is not None else 0.0,
                },
                "font_size_pt": font_size,
            })
    return out


def inspect(pptx_path: Path, page_idx: int) -> dict:
    if not pptx_path.exists():
        print(f"PPTX_NOT_FOUND: {pptx_path}", file=sys.stderr)
        sys.exit(1)
    try:
        p = Presentation(str(pptx_path))
    except Exception as e:
        print(f"PPTX_PARSE_ERROR: {e}", file=sys.stderr)
        sys.exit(2)

    if page_idx < 0 or page_idx >= len(p.slides):
        print(f"PAGE_OUT_OF_RANGE: page_idx={page_idx}, total={len(p.slides)}", file=sys.stderr)
        sys.exit(1)

    slide = p.slides[page_idx]
    leaves = walk_shapes(slide.shapes)
    # Sort by geometry: top first, then left
    leaves.sort(key=lambda s: (s["bbox"]["top"], s["bbox"]["left"]))

    slots = []
    for leaf in leaves:
        slots.append({
            "id": "?",
            "tree_path": leaf["tree_path"],
            "raw_text_sample": leaf["raw_text_sample"],
            "bbox": leaf["bbox"],
            "font_size_pt": leaf["font_size_pt"],
            "capacity_chars": "?",
            "text_color_override": None,
        })

    return {
        "template_page_index": page_idx,
        "layout_class": "?",
        "slots": slots,
    }


def main():
    if len(sys.argv) != 3:
        print("Usage: inspect_placeholders.py <pptx_path> <page_idx>", file=sys.stderr)
        sys.exit(1)
    result = inspect(Path(sys.argv[1]), int(sys.argv[2]))
    print(yaml.safe_dump(result, sort_keys=False, allow_unicode=True))


if __name__ == "__main__":
    main()
