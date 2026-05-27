"""One-shot script: generate sample.pptx fixture. Run once, commit the .pptx."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

p = Presentation()
# Page 1: cover · title + subtitle
slide1 = p.slides.add_slide(p.slide_layouts[5])  # blank
title = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
title.text_frame.text = "深蓝主标题"
title.text_frame.paragraphs[0].runs[0].font.size = Pt(44)
subtitle = slide1.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
subtitle.text_frame.text = "副标题占位"

# Page 2: cards · 3 textbox (simulating 3-cols cards)
slide2 = p.slides.add_slide(p.slide_layouts[5])
for i in range(3):
    tb = slide2.shapes.add_textbox(Inches(0.5 + i * 3), Inches(2), Inches(2.5), Inches(2))
    tb.text_frame.text = f"卡片 {i + 1}"

out = Path(__file__).parent / "sample.pptx"
p.save(out)
print(f"wrote {out}")
