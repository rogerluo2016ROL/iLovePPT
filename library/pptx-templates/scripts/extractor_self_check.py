"""Extractor Step 3.3 self-check · exit code 驱动 · 替代内嵌 bash.

Exit codes:
  0 = 全过
  1 = 字段缺失 / enum / 格式问题
  2 = placeholder_map tree_path 不能 resolve(此 task 暂不实现,Task 2 补)
  3 = YAML 语法错
  4 = 模板目录不存在
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

import yaml

ALLOWED_LAYOUTS = {
    "cover", "toc", "section_divider", "summary", "closing", "quote",
    "single_focus", "cards", "bullet_list", "data",
    "timeline", "pyramid", "venn", "radial", "process_flow", "quadrant", "comparison",
    "other",
}

REQUIRED_TEMPLATE_FIELDS = [
    "id", "name", "category", "content_intent", "when_to_use",
    "keywords", "recommended_for", "visual_signature",
    "status", "provenance", "extraction",
]

REQUIRED_PAGE_FIELDS = [
    "id", "name", "layout_type", "content_intent", "when_to_use",
    "keywords", "native_elements", "status", "confidence",
]

ID_RE = re.compile(r"^[a-z0-9_-]+__\d{2}-[a-z_]+$")


def _resolve_tree_path(shapes, tree_path: str) -> bool:
    """tree_path like '3' / '3.16' / '3.16.0' · walks shape tree · True if resolvable."""
    if not tree_path:
        return False
    parts = tree_path.split(".")
    current = list(shapes)
    for i, part in enumerate(parts):
        try:
            idx = int(part)
        except ValueError:
            return False
        if idx >= len(current):
            return False
        node = current[idx]
        if i == len(parts) - 1:
            return True
        try:
            current = list(node.shapes)
        except AttributeError:
            return False
    return True


def load_yaml(path: Path) -> tuple[dict | None, str | None]:
    """返回 (data, error_msg). data is None 时 error_msg 非空."""
    try:
        with open(path, encoding="utf-8") as f:
            data = yaml.safe_load(f)
    except yaml.YAMLError as e:
        return None, str(e).split("\n")[0]
    except FileNotFoundError:
        return None, f"file not found: {path}"
    if data is None:
        return None, "file is empty or null"
    if not isinstance(data, dict):
        return None, f"YAML_NOT_A_DICT: top-level is {type(data).__name__}, expected dict"
    return data, None


def check(name: str, items_root: Path) -> int:
    tpl_dir = items_root / name
    if not tpl_dir.is_dir():
        print(f"TEMPLATE_NOT_FOUND: {tpl_dir}")
        return 4

    errors: list[str] = []
    syntax_errors: list[str] = []

    # 0. YAML 语法(模板级 + 所有页)
    tpl_meta_path = tpl_dir / "meta.yaml.draft"
    if not tpl_meta_path.exists():
        tpl_meta_path = tpl_dir / "meta.yaml"
    tpl_meta, err = load_yaml(tpl_meta_path)
    if err:
        syntax_errors.append(f"YAML_SYNTAX_ERROR: {tpl_meta_path}: {err}")

    page_paths = sorted(tpl_dir.glob("pages/*/meta.yaml.draft"))
    if not page_paths:
        page_paths = sorted(tpl_dir.glob("pages/*/meta.yaml"))
    page_metas: list[tuple[Path, dict]] = []
    for p in page_paths:
        data, err = load_yaml(p)
        if err:
            syntax_errors.append(f"YAML_SYNTAX_ERROR: {p}: {err}")
        elif data is not None:
            page_metas.append((p, data))

    if syntax_errors:
        for e in syntax_errors:
            print(e)
        return 3

    # 1. 模板级必填字段
    if tpl_meta:
        for f in REQUIRED_TEMPLATE_FIELDS:
            if f not in tpl_meta:
                errors.append(f"MISSING_TEMPLATE_FIELD: {tpl_meta_path}: {f}")

    # 2. 页级必填字段
    for p, data in page_metas:
        for f in REQUIRED_PAGE_FIELDS:
            if f not in data:
                errors.append(f"MISSING_PAGE_FIELD: {p}: {f}")

    # 3. layout_type enum (None/"" treated as invalid, not skipped)
    for p, data in page_metas:
        lt = data.get("layout_type")
        if not isinstance(lt, str) or not lt:
            errors.append(f"LAYOUT_TYPE_INVALID: {p}: layout_type={lt!r} (must be non-empty string)")
        elif lt not in ALLOWED_LAYOUTS:
            errors.append(f"ENUM_VIOLATION: {p}: layout_type={lt}")

    # 4. id 格式 + 唯一性
    ids = []
    for p, data in page_metas:
        page_id = data.get("id", "")
        if not ID_RE.match(page_id):
            errors.append(f"ID_FORMAT_INVALID: {p}: id={page_id!r} (expected <name>__<NN-slug>)")
        ids.append(page_id)
    if len(set(ids)) != len(ids):
        errors.append(f"ID_DUPLICATE: {len(ids)} pages but {len(set(ids))} unique ids")

    # 5. confidence 是数字
    for p, data in page_metas:
        c = data.get("confidence")
        if not isinstance(c, (int, float)):
            errors.append(f"CONFIDENCE_NOT_NUMERIC: {p}: confidence={c!r}")
        elif not 0.0 <= c <= 1.0:
            errors.append(f"CONFIDENCE_OUT_OF_RANGE: {p}: confidence={c}")

    # 6. provenance.embedding_dim == 1152
    if tpl_meta:
        prov = tpl_meta.get("provenance")
        if not isinstance(prov, dict):
            prov = {}
        if prov.get("embedding_dim") != 1152:
            errors.append(f"EMBEDDING_DIM_WRONG: {tpl_meta_path}: got {prov.get('embedding_dim')}, expected 1152")

    # 7. extraction 算式自洽
    if tpl_meta:
        ext = tpl_meta.get("extraction")
        if not isinstance(ext, dict):
            ext = {}
        d = ext.get("declared_pages")
        r = ext.get("rendered_pages")
        disc = ext.get("discrepancy")
        if d is not None and r is not None and disc is not None:
            if not all(isinstance(x, int) and not isinstance(x, bool) for x in (d, r, disc)):
                errors.append(f"EXTRACTION_TYPE_INVALID: {tpl_meta_path}: declared/rendered/discrepancy must be int, got ({type(d).__name__}, {type(r).__name__}, {type(disc).__name__})")
            elif d - r != disc:
                errors.append(f"EXTRACTION_MATH_INCONSISTENT: declared={d} rendered={r} discrepancy={disc}")

    # 8. template_name 跟父目录名一致
    for p, data in page_metas:
        tn = data.get("template_name")
        if tn and tn != name:
            errors.append(f"TEMPLATE_NAME_MISMATCH: {p}: template_name={tn} != {name}")

    # 9. placeholder_map tree_path resolves (if pmap.draft exists AND source .pptx exists)
    pmap_errors: list[str] = []
    source_pptx = items_root.parent / "_source" / f"{name}.pptx"
    if source_pptx.exists():
        try:
            from pptx import Presentation as _Pres
            pres = _Pres(str(source_pptx))
            for p in tpl_dir.glob("pages/*/placeholder_map.yaml.draft"):
                pmap, err = load_yaml(p)
                if err or not pmap:
                    continue
                idx = pmap.get("template_page_index")
                if idx is None or not isinstance(idx, int) or idx < 0 or idx >= len(pres.slides):
                    pmap_errors.append(f"PMAP_PAGE_INDEX_INVALID: {p}: template_page_index={idx}")
                    continue
                slide = pres.slides[idx]
                for slot in pmap.get("slots", []):
                    tp = slot.get("tree_path", "")
                    if not _resolve_tree_path(slide.shapes, tp):
                        pmap_errors.append(f"PMAP_TREE_PATH_UNRESOLVABLE: {p}: tree_path={tp!r}")
        except Exception as e:
            pmap_errors.append(f"PMAP_CHECK_FAILED: {e}")

    # If ONLY pmap errors (no other errors), exit 2; otherwise normal exit-1 path
    if pmap_errors and not errors:
        for e in pmap_errors:
            print(e)
        return 2
    # If both pmap and other errors, merge into general errors → exit 1
    errors.extend(pmap_errors)

    if errors:
        for e in errors:
            print(e)
        return 1

    print(f"OK: {name} · {len(page_metas)} pages · all checks passed")
    return 0


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("name")
    ap.add_argument("--items-root", default="library/pptx-templates/items",
                    help="items/ 目录(默认 library/pptx-templates/items)")
    args = ap.parse_args()
    sys.exit(check(args.name, Path(args.items_root)))


if __name__ == "__main__":
    main()
